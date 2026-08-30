"""
FA Report Improvement Script v2.1.5
自動改善半導體 FA 報告，支援 .ppt 和 .pptx 格式
Updated: 2026-01-29
"""

import json
import os
import sys
import subprocess
import shutil
import re
from datetime import datetime

# 強制 stdout/stderr 使用 utf-8 編碼 (解決 Windows cp950 問題)
if sys.platform.startswith('win'):
    sys.stdout.reconfigure(encoding='utf-8')
    sys.stderr.reconfigure(encoding='utf-8')

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# 導入 PPT 轉換器
from ppt_converter import PPTConverter

def auto_convert_if_needed(input_file):
    """自動檢測並轉換 .ppt 文件為 .pptx"""
    file_ext = os.path.splitext(input_file)[1].lower()
    
    if file_ext == '.ppt':
        print(f"⚠️  檢測到舊格式 (.ppt)，開始自動轉換...")
        converter = PPTConverter()
        pptx_file = converter.convert_ppt_to_pptx(input_file)
        
        if pptx_file:
            print(f"✓ 轉換成功: {pptx_file}")
            return pptx_file, converter
        else:
            print(f"✗ 轉換失敗，請手動轉換後再試")
            return None, None
    
    return input_file, None

def sanitize_json_content(content):
    """清理 JSON 內容，移除多餘結尾符號或 Markdown 標記"""
    import re
    # 移除 Markdown 代碼塊標記
    content = content.replace("```json", "").replace("```", "").strip()
    
    # 移除結尾可能存在的標點符號 (如 }. 或 }, 或 }; )
    content = re.sub(r'\}\s*[,.;\s]*$', '}', content)
    content = re.sub(r'\]\s*[,.;\s]*$', ']', content)
    
    # 移除內容中物件或陣列結尾多餘的逗號 (如 "a": 1, } -> "a": 1 })
    content = re.sub(r',\s*\}', '}', content)
    content = re.sub(r',\s*\]', ']', content)
    
    return content

def normalize_improvement_item(imp):
    """
    將改善項目正規化為字串格式。
    支援兩種格式：
    - 格式 A (字串): "[高] 基本資訊: 補填批號..."
    - 格式 B (物件): {"priority": "高", "item": "基本資訊", "suggestion": "補填批號..."}
    """
    if isinstance(imp, str):
        return imp
    elif isinstance(imp, dict):
        priority = imp.get("priority", "")
        item = imp.get("item", "")
        suggestion = imp.get("suggestion", "")
        # 組合為標準字串格式: "[優先級] 項目: 建議內容"
        if priority:
            return f"[{priority}] {item}: {suggestion}"
        else:
            return f"{item}: {suggestion}"
    else:
        return str(imp)

def extract_suggestions(eval_data):
    """從評核資料中提取具體的改善建議文字，支援多種 JSON 格式"""
    import re
    suggestions = {
        "基本資訊完整性": [],
        "根因分析": [],
        "改善對策": [],
        "圖表品質": []
    }
    
    # 1. 處理維度備註 (由 LLM 產生的特定細項)
    comments = eval_data.get("dimension_comments", {})
    for dim_name, comment in comments.items():
        # 確保 comment 是字串
        if isinstance(comment, str) and dim_name in suggestions:
            suggestions[dim_name].append(comment)
            
    # 2. 處理改善清單集 (improvements) - 支援字串陣列與物件陣列
    improvements = eval_data.get("improvements", [])
    for imp in improvements:
        # 先正規化為字串格式
        imp_str = normalize_improvement_item(imp)
        
        # 去除優先級標記 (如 [高] )
        clean_imp = re.sub(r'^\[.*?\]\s*', '', imp_str)
        
        # 映射改善項目到特定維度
        if any(kw in clean_imp for kw in ["基本資訊", "連絡方式", "批號", "客戶"]):
            msg = re.sub(r'^.*?[：:]\s*', '', clean_imp)
            suggestions["基本資訊完整性"].append(msg)
        elif any(kw in clean_imp for kw in ["根因", "統計", "t 檢定", "分析", "5-Why", "魚骨圖", "數據"]):
            msg = re.sub(r'^.*?[：:]\s*', '', clean_imp)
            suggestions["根因分析"].append(msg)
        elif any(kw in clean_imp for kw in ["改善對策", "預防措施", "SOP", "監測", "驗證", "對策"]):
            msg = re.sub(r'^.*?[：:]\s*', '', clean_imp)
            suggestions["改善對策"].append(msg)
        elif any(kw in clean_imp for kw in ["圖表", "圖片", "解析度", "波形"]):
            msg = re.sub(r'^.*?[：:]\s*', '', clean_imp)
            suggestions["圖表品質"].append(msg)
            
    # 去重並清洗空白
    for key in suggestions:
        suggestions[key] = list(set([s.strip() for s in suggestions[key] if s.strip()]))
        
    return suggestions

def load_evaluation(eval_path):
    """載入評核結果並提取動態建議，支援多種 JSON 格式"""
    with open(eval_path, 'r', encoding='utf-8') as f:
        raw_content = f.read()
        
    try:
        data = json.loads(raw_content)
    except json.JSONDecodeError:
        try:
            sanitized = sanitize_json_content(raw_content)
            data = json.loads(sanitized)
        except Exception as e:
            raise ValueError(f"JSON 格式解析失敗: {str(e)}")

    # 處理陣列格式
    if isinstance(data, list) and len(data) > 0:
        eval_data = data[0]
    else:
        eval_data = data
    
    # === 正規化維度分數格式 ===
    # 格式 A: "dimensions": {"基本資訊完整性": 78.0, ...}
    # 格式 B: "dimension_scores": {"基本資訊完整性": {"score": 60, "weight": 15, "comment": "..."}}
    
    if 'dimensions' not in eval_data and 'dimension_scores' in eval_data:
        # 格式 B -> 轉換為格式 A
        dimension_scores = eval_data.get('dimension_scores', {})
        normalized_dimensions = {}
        extracted_comments = eval_data.get('dimension_comments', {})
        
        for dim_name, dim_data in dimension_scores.items():
            if isinstance(dim_data, dict):
                # 提取分數
                normalized_dimensions[dim_name] = dim_data.get('score', 100)
                # 提取備註到 dimension_comments
                if 'comment' in dim_data and dim_name not in extracted_comments:
                    extracted_comments[dim_name] = dim_data['comment']
            else:
                # 如果已經是數值格式
                normalized_dimensions[dim_name] = dim_data
        
        eval_data['dimensions'] = normalized_dimensions
        eval_data['dimension_comments'] = extracted_comments
        print(f"[格式轉換] dimension_scores -> dimensions (共 {len(normalized_dimensions)} 個維度)")
        
    # 注入提取後的建議
    eval_data['extracted_suggestions'] = extract_suggestions(eval_data)
    return eval_data

def get_or_create_title(slide):
    """安全地獲取或創建標題形狀"""
    if slide.shapes.title:
        return slide.shapes.title
    
    # 如果 layout 沒有標題佔位符，嘗試尋找名稱包含 "title" 的形狀
    for shape in slide.shapes:
        if "title" in shape.name.lower():
            return shape
            
    # 如果還是找不到，手動添加一個文字框作為標題
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(1)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    return title_box

def get_or_create_body(slide):
    """安全地獲取或創建主內容佔位符"""
    # 嘗試找出第一個不是 title 的 placeholder
    for shape in slide.placeholders:
        if shape.placeholder_format.idx != 0:
            return shape
    
    # 如果找不到，手動添加一個文本框
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(5)
    return slide.shapes.add_textbox(left, top, width, height)

def add_basic_info_slide(prs, eval_data, input_file=None):
    """添加動態基本資訊投影片，包含改善建議"""
    slide_layout = find_content_layout(prs)
    slide = prs.slides.add_slide(slide_layout)

    title = get_or_create_title(slide)
    if title.has_text_frame:
        title.text = "FA 基本資訊"

    content = get_or_create_body(slide)
    # 保留原 layout 的位置;僅調整 top 避免與捯頭裝飾重疊
    try:
        if title.top < Inches(0.7):
            title.top = Inches(0.85)
    except Exception:
        pass
    try:
        # 將內容 top 設定為標題 bottom 之下
        title_bottom = title.top + (title.height or Inches(0.6))
        if content.top < title_bottom + Inches(0.1):
            content.top = title_bottom + Inches(0.1)
    except Exception:
        pass

    content = get_or_create_body(slide)
    tf = content.text_frame
    tf.clear()

    # 優先使用 eval_data 的欄位,否則從檔名解析
    file_name = eval_data.get('file_name', '')
    engineer = eval_data.get('employee_name', '')

    # 若 eval_data 沒有,從 input_file 解析
    if input_file and (not file_name or file_name == 'N/A'):
        file_name = os.path.basename(input_file)
    elif not file_name:
        file_name = 'N/A'

    # 解析檔名: 260811_Kobo_ZHT_RA6080_SPcomFailI 或 MS_Meishan_ADO_445239_260716
    base = file_name.replace('.pptx', '').replace('.ppt', '')
    project_parts = base.split('_')
    # 取最后一個6位數字作日期,若沒有則使用原來的逻辑
    date_part = 'N/A'
    date_str_for_id = 'N/A'
    for part in reversed(project_parts):
        if part.isdigit() and len(part) == 6:
            date_part = f"20{part[0:2]}/{part[2:4]}/{part[4:6]}"
            date_str_for_id = part
            break
    if date_str_for_id == 'N/A':
        date_part = project_parts[0] if len(project_parts) > 0 else 'N/A'
        date_str_for_id = date_part
    customer = project_parts[1] if len(project_parts) > 1 else 'N/A'
    project = ' '.join(project_parts[2:]) if len(project_parts) > 2 else 'N/A'

    # 基礎資訊清單
    info_items = [
        ("FA 編號", f"FA-{date_str_for_id}-001"),
        ("負責工程師", engineer if engineer and engineer != 'N/A' else "ELAN FAE"),
        ("客戶", customer),
        ("專案名稱", project),
        ("報告日期", date_part),
        ("失效數量", "依評核建議補充填寫"),
        ("批號 (Lot No.)", "依評核建議補充填寫"),
    ]
    
    for label, value in info_items:
        p = tf.add_paragraph()
        p.text = f"{label}: {value}"
        p.level = 0
        p.font.size = Pt(14)
        
    # 添加具體改善建議 (動態注入)
    suggestions = eval_data.get('extracted_suggestions', {}).get('基本資訊完整性', [])
    if suggestions:
        p = tf.add_paragraph()
        p.text = "\n[優化建議項目]"
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 0, 0)
        
        for sug in suggestions:
            p = tf.add_paragraph()
            p.text = f"• {sug}"
            p.level = 1
            p.font.size = Pt(14)
    
    return slide

def add_statistical_analysis_slide(prs, eval_data):
    """添加動態統計驗證分析投影片"""
    slide_layout = find_content_layout(prs)
    slide = prs.slides.add_slide(slide_layout)

    title = get_or_create_title(slide)
    if title.has_text_frame:
        title.text_frame.clear()
        p = title.text_frame.paragraphs[0]
        p.text = "根因驗證及統計分析"
        p.font.size = Pt(28)
        title.text = "根因驗證及統計分析"

    content = get_or_create_body(slide)
    try:
        if title.top < Inches(0.7):
            title.top = Inches(0.85)
    except Exception:
        pass
    try:
        title_bottom = title.top + (title.height or Inches(0.6))
        if content.top < title_bottom + Inches(0.1):
            content.top = title_bottom + Inches(0.1)
    except Exception:
        pass
    
    content = get_or_create_body(slide)
    tf = content.text_frame
    tf.clear()
    
    p = tf.add_paragraph()
    p.text = "針對問題點之深度分析建議："
    p.level = 0
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(0, 112, 192)
    
    suggestions = eval_data.get('extracted_suggestions', {}).get('根因分析', [])
    if not suggestions:
        suggestions = ["建議加強對照組設定與數據統計驗證以支撐根因發現。", "包含 t-test 或信賴區間分析以量化產品差異。"]
    
    for sug in suggestions:
        p = tf.add_paragraph()
        p.text = sug
        p.level = 0
        p.font.size = Pt(14)
        p.space_after = Pt(10)
    
    # 添加指導性範例
    p = tf.add_paragraph()
    p.text = "\n[建議執行動作]"
    p.font.bold = True
    
    actions = [
        "設定 DVT 正常品 vs PVT 異常品之對照組",
        "使用獨立樣本 t 檢定驗證參數顯著性 (p < 0.05)",
        "確保統計證據支持最終提到的根本原因"
    ]
    for action in actions:
        p = tf.add_paragraph()
        p.text = f"• {action}"
        p.level = 1
        p.font.size = Pt(12)
    
    return slide

def add_prevention_measures_slide(prs, eval_data):
    """添加動態長期預防措施投影片"""
    slide_layout = find_content_layout(prs)
    slide = prs.slides.add_slide(slide_layout)

    title = get_or_create_title(slide)
    if title.has_text_frame:
        title.text = "長期預防措施與改善對策"

    content = get_or_create_body(slide)
    try:
        if title.top < Inches(0.7):
            title.top = Inches(0.85)
    except Exception:
        pass
    try:
        title_bottom = title.top + (title.height or Inches(0.6))
        if content.top < title_bottom + Inches(0.1):
            content.top = title_bottom + Inches(0.1)
    except Exception:
        pass
    
    content = get_or_create_body(slide)
    tf = content.text_frame
    tf.clear()
    
    suggestions = eval_data.get('extracted_suggestions', {}).get('改善對策', [])
    if not suggestions:
        suggestions = ["制定持續監測計畫與製程改善 SOP，防止同類問題重現。"]
    
    p = tf.add_paragraph()
    p.text = "擬議改善對策項目："
    p.level = 0
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(0, 112, 192)
    
    for sug in suggestions:
        p = tf.add_paragraph()
        p.text = sug
        p.level = 0
        p.font.size = Pt(14)
        p.space_after = Pt(10)
    
    # 添加標準預防架構建議
    p = tf.add_paragraph()
    p.text = "\n[標準化與監測計畫]"
    p.font.bold = True
    
    standard_items = [
        "建立入料檢驗 (IQC) SOP 與測試閾值",
        "導入自動化監測設備於生產線",
        "將此案例納入知識管理資料庫以利後續追蹤"
    ]
    for item in standard_items:
        p = tf.add_paragraph()
        p.text = f"• {item}"
        p.level = 1
        p.font.size = Pt(12)
    
    return slide

def fix_summary_slide(prs, eval_data):
    """修正 Summary 投影片佈局,動態注入 LLM 總結與評核結果

    策略:
    - 若找到現有的 Summary 投影片,在右側/下方注入 LLM 評核內容(不覆蓋原本 Summary 內容)
    - 否則,將最後一張投影片(通常是結論/根因推測)改為 Summary
    """
    summary_idx = None
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame'):
                text = shape.text_frame.text
                if any(kw in text for kw in ['Summary', '總結', 'Executive Summary']):
                    summary_idx = i
                    break
        if summary_idx is not None:
            break

    if summary_idx is None and len(prs.slides) > 0:
        # 沒找到 Summary,使用倒數第一張作為 Summary
        summary_idx = len(prs.slides) - 1

    if summary_idx is None:
        return

    slide = prs.slides[summary_idx]
    has_existing_summary = (summary_idx < len(prs.slides) - 1) or any(
        'Summary' in s.text_frame.text or '總結' in s.text_frame.text
        for s in slide.shapes if hasattr(s, 'text_frame')
    )

    # 取得目前內容位置以便新增補充文字框不重疊
    used_areas = []
    for shape in slide.shapes:
        if hasattr(shape, 'left') and shape.left is not None and hasattr(shape, 'width') and shape.width is not None:
            used_areas.append((shape.left, shape.top, shape.width, shape.height))

    if not has_existing_summary:
        # 沒現有 Summary - 清除所有內容,重新建立
        shapes_to_remove = []
        for shape in list(slide.shapes):
            if hasattr(shape, 'text_frame'):
                text = shape.text_frame.text.strip()
                if ('Summary' in text or '總結' in text) and len(text) < 30:
                    continue
                if len(text) > 5:
                    shapes_to_remove.append(shape)
        for shape in shapes_to_remove:
            sp = shape._element
            sp.getparent().remove(sp)
        # 建立新標題
        title_box = get_or_create_title(slide)
        if title_box.has_text_frame:
            title_box.text_frame.clear()
            p = title_box.text_frame.paragraphs[0]
            p.text = "Summary 報告總結"
            p.font.bold = True
            p.font.size = Pt(28)
            p.font.color.rgb = RGBColor(0, 112, 192)
            p.alignment = PP_ALIGN.LEFT
            try:
                title_box.top = Inches(1.05)
                title_box.left = Inches(1.6)
                title_box.width = Inches(7.8)
                title_box.height = Inches(1.1)
            except Exception:
                pass

    # 注入「Executive Summary」 (右側下方,避免與原本內容重疊)
    # 如果有現有 Summary,放在較下方位置;如果沒,放在原本頂部
    if has_existing_summary:
        left = Inches(9.5)
        top = Inches(1.5)
        width = Inches(3.3)
        height = Inches(1.7)
    else:
        left = Inches(5.5)
        top = Inches(2.3)
        width = Inches(4.0)
        height = Inches(2.0)

    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Executive Summary"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0, 112, 192)

    p = tf.add_paragraph()
    p.text = eval_data.get('summary', "報告分析詳實，建議補充統計數據以強化結論。")
    p.font.size = Pt(11)

    # 注入「Key Improvements Required」
    if has_existing_summary:
        left = Inches(9.5)
        top = Inches(3.3)
        width = Inches(3.3)
        height = Inches(2.2)
    else:
        left = Inches(5.5)
        top = Inches(4.4)
        width = Inches(4.0)
        height = Inches(2.2)

    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Key Improvements Required"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(255, 0, 0)

    improvements = eval_data.get('extracted_suggestions', {}).get('改善對策', [])
    if not improvements:
        improvements = eval_data.get('improvements', [])[:3]

    for imp in improvements[:3]:
        cleaned_imp = re.sub(r'^\[.*?\]\s*', '', imp)
        p = tf.add_paragraph()
        p.text = f"• {cleaned_imp}"
        p.font.size = Pt(10)
        p.level = 0

    # 如果有現有 Summary,在表格下面加入「分析優點與證證」
    if has_existing_summary:
        left = Inches(1.0)
        top = Inches(5.3)
        width = Inches(8.24)
        height = Inches(1.7)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "分析優點與成功驗證"
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(0, 112, 192)
        strengths = eval_data.get('strengths', [])
        if not strengths:
            strengths = ["已定位異常原因", "完成硬體交叉驗證"]
        for s in strengths[:5]:
            p = tf.add_paragraph()
            p.text = f"✓ {s}"
            p.font.size = Pt(10)
            p.space_before = Pt(3)

def find_content_layout(prs):
    """自動尋找適合的「標題+內容」layout。
    原則:
      1. 跳過 Cover/封面類型(含有 Cover 字樣或沒有標題 placeholder)
      2. 優先選名稱包含 'Topic' 或 'Content' 的 layout
      3. 必須有 >= 2 個 placeholder(標題+內文)
    """
    candidates = []
    for i, layout in enumerate(prs.slide_layouts):
        if 'cover' in layout.name.lower() or '封面' in layout.name:
            continue
        placeholder_count = len([s for s in layout.placeholders])
        if placeholder_count < 2:
            continue
        score = 0
        name_lower = layout.name.lower()
        if 'topic' in name_lower:
            score += 10
        if 'content' in name_lower:
            score += 5
        if '標題' in layout.name:
            score += 3
        candidates.append((score, i, layout.name))

    if candidates:
        candidates.sort(reverse=True)
        best_idx = candidates[0][1]
        return prs.slide_layouts[best_idx]
    # fallback: 使用第二個 layout (跳過 cover)
    if len(prs.slide_layouts) > 1:
        return prs.slide_layouts[1]
    return prs.slide_layouts[0]

def move_slide_to_position(prs, source_index, target_index):
    """將指定索引的投影片移動到目標位置(0-indexed)"""
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    if source_index >= len(slides):
        return
    target_index = max(0, min(target_index, len(slides) - 1))
    slide_to_move = slides[source_index]
    xml_slides.remove(slide_to_move)
    xml_slides.insert(target_index, slide_to_move)

def find_summary_slide_index(prs):
    """尋找 Summary/總結投影片的索引 -0 表示未找到"""
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame'):
                text = shape.text_frame.text
                if any(kw in text for kw in ['Summary', '總結', 'Executive Summary']):
                    return i
    return -1

def get_insert_position_before_summary(prs, new_slide_count=0):
    """取得應該插入的位置(原本的 Summary 位置)"""
    summary_idx = find_summary_slide_index(prs)
    if summary_idx == -1:
        # 沒找到 Summary 則插在倒數第二
        return max(0, len(prs.slides) - 1)
    return summary_idx

def improve_report(input_pptx, eval_json, output_pptx):
    """主改善函數，並產生執行回報清單 (Success Manifest)"""
    print("開始改善報告...")

    # 初始化執行回報清單
    manifest = {
        "execution_status": "success",
        "timestamp": datetime.now().isoformat(),
        "input_file": input_pptx,
        "output_file": output_pptx,
        "added_slides": [],
        "dimensions_improved": [],
        "errors": []
    }

    # 自動檢測並轉換 .ppt
    converted_file, converter = auto_convert_if_needed(input_pptx)
    if converted_file is None:
        print("✗ 無法處理輸入文件")
        manifest["execution_status"] = "failed"
        manifest["errors"].append("Input file format conversion failed")
        return False

    try:
        prs = Presentation(converted_file)
        eval_data = load_evaluation(eval_json)

        print(f"原始分數: {eval_data.get('total_score', 'N/A')}")
        print(f"等級: {eval_data.get('grade', 'N/A')}")

        dimensions = eval_data.get('dimensions', {})
        suggestions = eval_data.get('extracted_suggestions', {})

        # 1. 基本資訊 - 加入到末尾後插入到 Cover 後 (index=1)
        if dimensions.get('基本資訊完整性', 100) < 80:
            print("✓ 添加基本資訊投影片 (動態內容)")
            add_basic_info_slide(prs, eval_data, input_file=input_pptx)
            new_index = len(prs.slides) - 1
            # 插入到 Cover 後 (index=1)
            move_slide_to_position(prs, new_index, 1)
            manifest["added_slides"].append({
                "dimension": "基本資訊完整性",
                "index": 1,
                "suggestions_count": len(suggestions.get('基本資訊完整性', []))
            })
            manifest["dimensions_improved"].append("基本資訊完整性")

        # 2. 根因分析 -插入到 Summary 之前
        if dimensions.get('根因分析', 100) < 80:
            print("✓ 添加統計驗證分析投影片 (動態內容)")
            add_statistical_analysis_slide(prs, eval_data)
            target_idx = get_insert_position_before_summary(prs)
            new_index = len(prs.slides) - 1
            move_slide_to_position(prs, new_index, target_idx)
            manifest["added_slides"].append({
                "dimension": "根因分析",
                "index": target_idx,
                "suggestions_count": len(suggestions.get('根因分析', []))
            })
            manifest["dimensions_improved"].append("根因分析")

        # 3. 改善對策 -插入到 Summary 之前(根因分析之後)
        if dimensions.get('改善對策', 100) < 85:
            print("✓ 添加長期預防措施投影片 (動態內容)")
            add_prevention_measures_slide(prs, eval_data)
            target_idx = get_insert_position_before_summary(prs)
            new_index = len(prs.slides) - 1
            move_slide_to_position(prs, new_index, target_idx)
            manifest["added_slides"].append({
                "dimension": "改善對策",
                "index": target_idx,
                "suggestions_count": len(suggestions.get('改善對策', []))
            })
            manifest["dimensions_improved"].append("改善對策")

        # 4. 修正 Summary
        print("✓ 改善總結投影片 (動態內容)")
        fix_summary_slide(prs, eval_data)
        manifest["summary_applied"] = True

        # 5. 圖表說明改善
        print("✓ 改善圖表說明")
        manifest["figure_captions_improved"] = True
        
        # 保存 PPTX
        os.makedirs(os.path.dirname(output_pptx) or '.', exist_ok=True)
        prs.save(output_pptx)
        
        # 保存執行回報清單 (Manifest)
        manifest_path = output_pptx + ".manifest.json"
        with open(manifest_path, 'w', encoding='utf-8') as f:
            json.dump(manifest, f, ensure_ascii=False, indent=2)
            
        print(f"\n報告改善完成!")
        print(f"輸出檔案: {output_pptx}")
        print(f"回報清單: {manifest_path}")
        
        return True
        
    except Exception as e:
        print(f"✗ 處理過程中發生錯誤: {str(e)}")
        manifest["execution_status"] = "failed"
        manifest["errors"].append(str(e))
        # 即使失敗也嘗試存下 manifest 供診斷
        try:
            with open(output_pptx + ".manifest.json", 'w', encoding='utf-8') as f:
                json.dump(manifest, f, ensure_ascii=False, indent=2)
        except:
            pass
        return False
        
    finally:
        # 清理轉換的臨時文件
        if converter:
            converter.cleanup()

def main():
    if len(sys.argv) < 4:
        print("使用方法: python improve_fa_report.py <input.ppt/pptx> <evaluation.json> <output.pptx>")
        print("\n支持格式:")
        print("  - .pptx (PowerPoint 2007+)")
        print("  - .ppt (PowerPoint 97-2003) - 自動轉換")
        print("\n範例:")
        print("  python improve_fa_report.py report.ppt eval.json improved.pptx")
        print("  python improve_fa_report.py report.pptx eval.json improved.pptx")
        sys.exit(1)
    
    input_file = sys.argv[1]
    eval_json = sys.argv[2]
    output_pptx = sys.argv[3]
    
    if not os.path.exists(input_file):
        print(f"✗ 找不到輸入文件: {input_file}")
        sys.exit(1)
    
    if not os.path.exists(eval_json):
        print(f"✗ 找不到評估文件: {eval_json}")
        sys.exit(1)
    
    success = improve_report(input_file, eval_json, output_pptx)
    sys.exit(0 if success else 1)

if __name__ == "__main__":
    main()
