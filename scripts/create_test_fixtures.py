#!/usr/bin/env python3
"""建立 CI 測試 fixtures

供 GitHub Actions workflow 使用,在執行 pytest 前動態建立:
- report/test_sample.pptx(測試 pptx)
- report/test_eval.json(測試評估 JSON)
- report/fa_report_test.txt(測試評估 TXT,符合 fa_report_analyzer 格式)
"""

from __future__ import annotations

import json

from pptx import Presentation

from fa_improver.paths import SKILL_ROOT


def main() -> None:
    """建立所有測試 fixtures"""
    # v3.1.5(P1):原本 Path("report") 依賴 cwd,從別的目錄執行會寫錯地方
    report_dir = SKILL_ROOT / "report"
    report_dir.mkdir(exist_ok=True)

    # 1. 建立測試 pptx(有母片 + 2 個 layouts)
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Test FA Report"
    prs.save(report_dir / "test_sample.pptx")

    # 2. 建立測試 eval JSON
    eval_data = {
        "total_score": 70.0,
        "grade": "C",
        "dimensions": {
            "基本資訊完整性": {"score": 60, "weight": 15, "comment": "需補充填寫"},
            "問題描述與定義": {"score": 70, "weight": 15, "comment": "OK"},
            "分析方法與流程": {"score": 75, "weight": 20, "comment": "OK"},
            "數據與證據支持": {"score": 70, "weight": 20, "comment": "OK"},
            "根因分析": {"score": 60, "weight": 20, "comment": "需加強"},
            "改善對策": {"score": 80, "weight": 10, "comment": "OK"},
        },
        "improvements": [{"item": "基本資訊", "suggestion": "補充填寫", "priority": "HIGH"}],
        "summary": "測試報告",
        "strengths": ["完整分析"],
    }
    (report_dir / "test_eval.json").write_text(
        json.dumps(eval_data, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )

    # 3. 建立測試 eval TXT(模擬 fa_report_analyzer 文字輸出)
    # 格式需符合 EvaluationParser._parse_dimensions_from_txt:
    #   【維度名稱】
    #   得分: 60 / 100
    #   權重: 15%
    #   評語: ...
    txt_content = """FA 報告評估結果

報告名稱: test_sample.pptx
評估日期: 2026-08-31

總分: 70.0 / 100
等級: C

【基本資訊完整性】
得分: 60 / 100 (60.0%)
權重: 15%
評語: 需補充填寫

【問題描述與定義】
得分: 70 / 100 (70.0%)
權重: 15%
評語: OK

【分析方法與流程】
得分: 75 / 100 (75.0%)
權重: 20%
評語: OK

【數據與證據支持】
得分: 70 / 100 (70.0%)
權重: 20%
評語: OK

【根因分析】
得分: 60 / 100 (60.0%)
權重: 20%
評語: 需加強對照組設定

【改善對策】
得分: 80 / 100 (80.0%)
權重: 10%
評語: OK

改善建議:
[基本資訊] 補充填寫
[根因分析] 需加強對照組設定
"""
    (report_dir / "fa_report_test.txt").write_text(txt_content, encoding="utf-8")

    print("✅ Test fixtures created in report/")


if __name__ == "__main__":
    main()
