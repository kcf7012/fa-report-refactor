"""量測母片左上裝飾的實際 x 範圍 —— 用來決定 `TITLE_SAFE_LEFT_INCH`。

跨平台遷移計劃書 P4 第 1 步:「先量測,再定值」。稽核明確警告過
**不要為了讓既有測試通過而反推數值**,所以常數必須先有量測依據,
再回頭決定值,而不是先挑一個值再湊理由。

## 量測範圍為什麼是「母片 + 被選中的那一個 layout」

所有 improver 都經由 `find_content_layout(prs)` 取得 layout,
**一份簡報只會用到一個 layout**(見 `improvers/*.py` 的 `add_slide`)。
其餘 layout(Agenda、End、Chapter Idea …)上的裝飾永遠不會出現在
新投影片上,把它們算進來只會讓數字虛高。腳本仍會印出全 layout 的
掃描結果當作參考,但結論只採計實際會用到的那一個。

## 什麼算「擋得住 title」

三類分開統計,因為它們對「把 title 往右移」的反應完全不同:

- **可迴避的左側裝飾** —— 靠左緣、右緣未過投影片中線。往右移就能閃開,
  這類的**最大右緣**才是 `TITLE_SAFE_LEFT_INCH` 的下限依據。
- **滿版背景 / 橫幅** —— 寬度 >= 投影片 90%。往右移一寸也閃不開,
  它們本來就是 title 要疊在上面的設計元素,不列入下限計算(但會印出來,
  免得「排除了什麼」變成看不見的黑箱)。
- **其他** —— 落在中段的裝飾,單獨列出供人判斷。

同一份掃描分成 **title 帶**與 **body 帶**兩段。body margin 目前是
`TITLE_SAFE_LEFT_INCH - 0.2`,它該不該調取決於 body 帶的裝飾,
拿 title 帶的數字去論證 body 的邊界會得到不適用的結論。

## 用法

    uv run python scripts/measure_master_decoration.py            # 真實客戶檔 + 合成 fixture
    uv run python scripts/measure_master_decoration.py a.pptx b.pptx

不帶參數時透過 `fa_improver.paths.resolve_report_file()` 尋找真實客戶檔
(在**根倉庫**的 `report/`,不是技能包自己那個只有 test_sample.pptx 的)。
找不到時明確報告缺哪一份,不靜默降級 —— 靜默降級正是本專案的歷史病灶。
"""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation

from fa_improver.improvers._safe_shape import TITLE_SAFE_HEIGHT_INCH, TITLE_SAFE_TOP_INCH
from fa_improver.layout.selector import find_content_layout
from fa_improver.paths import SKILL_ROOT, resolve_report_file

EMU_PER_INCH = 914400

# 三份真實客戶報告(根倉庫 report/)
REAL_REPORTS = (
    "260811_Kobo_ZHT_RA6080_SPcomFailI.pptx",
    "MS_Meishan_ADO_445239_260716.pptx",
    "N160JCN-EEK project 1pcs NG sample analysis report 260810.pptx",
)

# 合成 fixture(CI 上唯一存在的素材,一併量測以確認測試素材有代表性)
SYNTHETIC_DIR = SKILL_ROOT / "tests" / "integration" / "_synthetic_fixtures"
SYNTHETIC_FIXTURES = (
    "synthetic_A_vertical.pptx",
    "synthetic_B_single_placeholder.pptx",
    "synthetic_C_decoration.pptx",
)

# title 帶的垂直範圍 —— 直接取自 _safe_shape 的常數,不另外寫死
TITLE_BAND_TOP = TITLE_SAFE_TOP_INCH
TITLE_BAND_BOTTOM = TITLE_SAFE_TOP_INCH + TITLE_SAFE_HEIGHT_INCH

# body 帶 —— 取自 get_or_create_body() 的 fallback 幾何(top=1.5、height=sh-2.0)
BODY_BAND_TOP = 1.5
BODY_BAND_BOTTOM_MARGIN = 0.5

# 分類門檻(比例,對投影片寬度取值,才能同時適用 10 in 與 13.33 in 母片)
FULL_BLEED_RATIO = 0.90  # 寬度佔比 >= 此值視為滿版背景/橫幅
LEFT_ANCHOR_RATIO = 0.25  # left 佔比 <= 此值視為靠左緣
ESCAPABLE_RATIO = 0.50  # 右緣佔比 < 此值才「往右移就閃得開」


def _inch(value) -> float | None:
    return None if value is None else value / EMU_PER_INCH


def _fmt(value: float | None) -> str:
    return "  n/a" if value is None else f"{value:5.2f}"


def _geometry(shape, source: str) -> dict:
    left, top = _inch(shape.left), _inch(shape.top)
    width, height = _inch(shape.width), _inch(shape.height)
    return {
        "source": source,
        "name": shape.name,
        "left": left,
        "top": top,
        "width": width,
        "height": height,
        "right": None if left is None or width is None else left + width,
        "bottom": None if top is None or height is None else top + height,
    }


def _decorations(shapes, source: str) -> list[dict]:
    """非 placeholder 的 shape 就是裝飾(Logo、色塊、直條、機密標示)。"""
    return [_geometry(s, source) for s in shapes if not s.is_placeholder]


def _in_band(row: dict, top: float, bottom: float) -> bool:
    if row["top"] is None or row["bottom"] is None:
        return True  # 沒有垂直資訊時保守地算它擋得到
    return row["top"] < bottom and row["bottom"] > top


def _classify(row: dict, slide_w: float) -> str:
    if row["left"] is None or row["right"] is None or row["width"] is None:
        return "unknown"
    if row["width"] >= slide_w * FULL_BLEED_RATIO:
        return "full_bleed"
    if row["left"] <= slide_w * LEFT_ANCHOR_RATIO and row["right"] < slide_w * ESCAPABLE_RATIO:
        return "left_escapable"
    return "other"


def measure(path: Path) -> dict:
    prs = Presentation(str(path))
    slide_w, slide_h = _inch(prs.slide_width), _inch(prs.slide_height)
    layout = find_content_layout(prs)

    used = _decorations(prs.slide_master.shapes, "master")
    used += _decorations(layout.shapes, f"layout:{layout.name}")

    all_layouts = list(used)
    for other in prs.slide_master.slide_layouts:
        if other is layout:
            continue
        all_layouts += _decorations(other.shapes, f"layout:{other.name}")

    bands = {
        "title": (TITLE_BAND_TOP, TITLE_BAND_BOTTOM),
        "body": (BODY_BAND_TOP, slide_h - BODY_BAND_BOTTOM_MARGIN),
    }

    def bucketize(rows, band):
        top, bottom = bands[band]
        buckets: dict[str, list[dict]] = {}
        for row in rows:
            if not _in_band(row, top, bottom):
                continue
            buckets.setdefault(_classify(row, slide_w), []).append(row)
        return buckets

    title_ph = None
    for ph in layout.placeholders:
        pf = ph.placeholder_format
        if pf.idx == 0 or str(pf.type).startswith(("TITLE", "CENTER_TITLE")):
            title_ph = _geometry(ph, f"layout:{layout.name}")
            break

    return {
        "path": path,
        "slide_w": slide_w,
        "slide_h": slide_h,
        "layout_name": layout.name,
        "layout_ph_count": len(list(layout.placeholders)),
        "title_ph": title_ph,
        "bands": bands,
        "used": {b: bucketize(used, b) for b in bands},
        "all": {b: bucketize(all_layouts, b) for b in bands},
    }


def _print_bucket(label: str, rows: list[dict]) -> None:
    print(f"  {label}({len(rows)}):")
    if not rows:
        print("      (無)")
        return
    seen = set()
    for row in sorted(rows, key=lambda r: r["right"] or 0):
        key = (row["source"], row["left"], row["top"], row["width"], row["height"])
        if key in seen:
            continue
        seen.add(key)
        print(
            f"      left={_fmt(row['left'])} top={_fmt(row['top'])} "
            f"w={_fmt(row['width'])} h={_fmt(row['height'])} "
            f"right={_fmt(row['right'])}  {row['source']} / {row['name']}"
        )


def _max_right(rows: list[dict]) -> float | None:
    return max((r["right"] for r in rows if r["right"] is not None), default=None)


def report(result: dict) -> dict[str, tuple[float | None, float | None]]:
    print(f"\n{'=' * 78}")
    print(f"檔案:{result['path'].name}")
    print(f"投影片:{result['slide_w']:.2f} x {result['slide_h']:.2f} in")
    print(
        f"find_content_layout() 選中:{result['layout_name']!r}"
        f"({result['layout_ph_count']} 個 placeholder)"
    )
    tp = result["title_ph"]
    if tp is None:
        print("  → 該 layout **沒有** title placeholder(idx=0 / TITLE 型別)")
    else:
        print(
            f"  → title placeholder:left={_fmt(tp['left'])} top={_fmt(tp['top'])} "
            f"w={_fmt(tp['width'])} h={_fmt(tp['height'])}  {tp['name']}"
        )

    limits = {}
    for band, (top, bottom) in result["bands"].items():
        print(f"\n[{band} 帶 {top:.2f}~{bottom:.2f} in](範圍:母片 + {result['layout_name']!r})")
        for key, label in (
            ("left_escapable", "可迴避的左側裝飾 ← 決定下限"),
            ("full_bleed", "滿版背景/橫幅(往右移閃不開,不列入下限)"),
            ("other", "其他"),
            ("unknown", "幾何不完整"),
        ):
            _print_bucket(label, result["used"][band].get(key, []))

        used_limit = _max_right(result["used"][band].get("left_escapable", []))
        all_limit = _max_right(result["all"][band].get("left_escapable", []))
        print(f"    實用範圍的可迴避裝飾最大右緣:{_fmt(used_limit)} in")
        print(f"    全 layout 掃描(參考值)      :{_fmt(all_limit)} in")
        limits[band] = (used_limit, all_limit)
    return limits


def main(argv: list[str]) -> int:
    if argv:
        paths = [Path(a) for a in argv]
        missing = [p for p in paths if not p.exists()]
        paths = [p for p in paths if p.exists()]
    else:
        resolved = {name: resolve_report_file(name) for name in REAL_REPORTS}
        missing = [Path(n) for n, p in resolved.items() if p is None]
        paths = [p for p in resolved.values() if p is not None]
        for name in SYNTHETIC_FIXTURES:
            candidate = SYNTHETIC_DIR / name
            (paths if candidate.exists() else missing).append(candidate)

    if missing:
        print("找不到以下檔案(量測結果不完整):", file=sys.stderr)
        for m in missing:
            print(f"  - {m}", file=sys.stderr)
    if not paths:
        print("沒有任何可量測的檔案。", file=sys.stderr)
        return 1

    results = []
    for path in paths:
        results.append((path.name, report(measure(path))))

    print(f"\n{'=' * 78}")
    print("彙總 —— 可迴避左側裝飾的最大右緣(in)")
    for band, hint in (
        ("title", "→ TITLE_SAFE_LEFT_INCH 的下限"),
        ("body", "→ body margin 的下限"),
    ):
        print(f"\n[{band} 帶] {hint}")
        print(f"  {'實用範圍':>8} {'全 layout':>10}  檔案")
        for name, limits in results:
            used_limit, all_limit = limits[band]
            print(f"  {_fmt(used_limit):>8} {_fmt(all_limit):>10}  {name}")
        used_all = [v for _, lim in results if (v := lim[band][0]) is not None]
        scan_all = [v for _, lim in results if (v := lim[band][1]) is not None]
        print(f"  實用範圍跨檔最大值:{max(used_all):.2f} in" if used_all else "  實用範圍:無")
        print(f"  全 layout 跨檔最大值:{max(scan_all):.2f} in" if scan_all else "  全 layout:無")
    return 0


if __name__ == "__main__":
    raise SystemExit(main(sys.argv[1:]))
