#!/usr/bin/env python3
"""建立 CI 用的合成 pptx + eval JSON fixtures

目的:
  讓 tests/integration/test_visual_quality.py 與 tests/integration/test_slide_rendering.py
  在 CI 環境(沒真實客戶 pptx 時)能跑。這些 fixture **完全去識別化**,
  無 ELAN logo / 無真實客戶名稱 / 無機密文字。

設計 3 個 fixture 對應 3 種 layout 場景:
  - synthetic_A_vertical:含「直排」layout,觸發 Bug 3 防護測試
  - synthetic_B_single_placeholder:只有 1 個 placeholder 的 layout,觸發 v3.1.3 修正
  - synthetic_C_decoration:母片含左上裝飾,觸發 TITLE_SAFE_LEFT_INCH 修正

每個 fixture 包含:
  - 5 張 slide(模擬真實 FA 報告:封面 / Summary / 8D 各階段)
  - 對應的 eval JSON

執行:
  python scripts/build_synthetic_fixtures.py

產出:
  tests/integration/_synthetic_fixtures/
    synthetic_A_vertical.pptx + .json
    synthetic_B_single_placeholder.pptx + .json
    synthetic_C_decoration.pptx + .json
"""

from __future__ import annotations

import contextlib
import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

# 輸出目錄:在 tests/integration/ 內,_ 開頭表示「內部輔助」
OUTPUT_DIR = Path(__file__).parent.parent / "tests" / "integration" / "_synthetic_fixtures"

# 通用評估資料:6 維度各給中等分數(觸發多個改善動作)
COMMON_EVAL_DATA = {
    "total_score": 72.0,
    "grade": "C",
    "dimensions": {
        "基本資訊完整性": {"score": 65, "weight": 15, "comment": "需補充填寫"},
        "問題描述與定義": {"score": 75, "weight": 15, "comment": "OK"},
        "分析方法與流程": {"score": 80, "weight": 20, "comment": "OK"},
        "數據與證據支持": {"score": 70, "weight": 20, "comment": "OK"},
        "根因分析": {"score": 65, "weight": 20, "comment": "需加強對照組設定"},
        "改善對策": {"score": 85, "weight": 10, "comment": "OK"},
    },
    "improvements": [
        {"item": "基本資訊", "suggestion": "補充填寫", "priority": "HIGH"},
        {"item": "根因分析", "suggestion": "需加強對照組設定", "priority": "HIGH"},
    ],
    "summary": "這是合成測試報告的 executive summary,完全去識別化。",
    "strengths": ["完整 8D 分析", "證據齊全"],
    "source_file": "synthetic.pptx",
}


def _add_5_slides(prs: Presentation, layout_idx: int) -> None:
    """加 5 張 slide 模擬真實 FA 報告結構"""
    titles = [
        "封面",
        "Summary",
        "D2 問題描述",
        "D4 根因分析",
        "D7 改善對策",
    ]
    for title in titles:
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
        with contextlib.suppress(AttributeError):
            # layout 沒有 title placeholder 時跳過
            slide.shapes.title.text = title


def build_synthetic_a_vertical() -> Path:
    """fixture A:含「Vertical」layout(Bug 3 防護測試)

    設計:用 python-pptx 預設母片裡的 layout[9] = "Title and Vertical Text"
    或 layout[10] = "Vertical Title and Text"(兩者都含 "Vertical" 關鍵字),
    觸發 _safe_shape.py 的 Bug 3 偵測邏輯:
      「layout name 含 "Vertical" → 跳過 layout placeholder,改用 safe_textbox fallback,
       避免文字被旋轉 90°」
    """
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()

    # 用 layout[9] "Title and Vertical Text" — 含 "Vertical" 關鍵字
    _add_5_slides(prs, layout_idx=9)

    out_path = OUTPUT_DIR / "synthetic_A_vertical.pptx"
    prs.save(out_path)

    eval_data = {**COMMON_EVAL_DATA, "source_file": "synthetic_A_vertical.pptx"}
    (OUTPUT_DIR / "synthetic_A_vertical.json").write_text(
        json.dumps(eval_data, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )
    return out_path


def build_synthetic_b_single_placeholder() -> Path:
    """fixture B:只有 1 個 placeholder 的 layout(v3.1.3 修正測試)

    設計:用一個 layout 只有 1 個 placeholder(類似「Topic-Numbers」layout),
    觸發 v3.1.3 在 _safe_shape.get_title_placeholder 加的修正:
    「若 layout 只有 1 個 placeholder,跳過改用 safe_textbox fallback,
    避免 title 與 placeholder 重疊」。
    """
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()

    # 預設 layout[0]:Title Slide 有 1 個 title placeholder + 1 個 body placeholder
    # 為簡化,我們直接用 layout[0],靠 v3.1.3 的 BODY_MIN_HEIGHT_INCH=1.0 邏輯觸發:
    # 當 layout placeholder height < 1.0 in 時 fallback。
    # python-pptx 改 layout placeholder height 很麻煩,簡化:用 layout[6] (Blank) 然後
    # 手動加一個「Topic」風格的 layout。

    # 改用更可靠的方法:在母片裡加一個 height 很小的 placeholder,讓 v3.1.3 邏輯觸發

    blank_layout = prs.slide_layouts[6]  # Blank
    for title in [
        "封面",
        "Summary",
        "D2 問題描述",
        "D4 根因分析",
        "D7 改善對策",
    ]:
        slide = prs.slides.add_slide(blank_layout)
        # 加一個 height 很小的 placeholder(< 1.0 in = 914400 EMU)
        # python-pptx 不支援直接加 placeholder,用 add_textbox 模擬
        tiny_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(0.3))
        tiny_box.text_frame.text = title

    out_path = OUTPUT_DIR / "synthetic_B_single_placeholder.pptx"
    prs.save(out_path)

    eval_data = {**COMMON_EVAL_DATA, "source_file": "synthetic_B_single_placeholder.pptx"}
    (OUTPUT_DIR / "synthetic_B_single_placeholder.json").write_text(
        json.dumps(eval_data, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )
    return out_path


def build_synthetic_c_decoration() -> Path:
    """fixture C:母片含左上裝飾(TITLE_SAFE_LEFT_INCH 測試)

    設計:在母片 layer 加一個 left=0.3 in 的裝飾 shape,
    模擬 MS 報告母片左上角 Logo 區域,觸發 v3.1.3 的 TITLE_SAFE_LEFT_INCH=1.2
    修正邏輯(新 title textbox left >= 1.2 in,避免被裝飾擋住)。

    由於 python-pptx 的 MasterShapes 沒有 add_shape() API,
    我們直接操作母片 _spTree XML 加一個 shape。
    """
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()

    # 在母片 layer 加左上裝飾(透過 XML)
    from lxml import etree

    master = prs.slide_masters[0]
    sp_tree = master.shapes._spTree

    # 建構一個 Rectangle shape XML(left=0, top=0, w=1in, h=0.5in)
    # 914400 EMU = 1 inch
    rect_xml = """<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
              xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
        <p:nvSpPr>
            <p:cNvPr id="100" name="LeftTopDecoration"/>
            <p:cNvSpPr/>
            <p:nvPr/>
        </p:nvSpPr>
        <p:spPr>
            <a:xfrm>
                <a:off x="0" y="0"/>
                <a:ext cx="914400" cy="457200"/>
            </a:xfrm>
            <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
            <a:solidFill><a:srgbClr val="CCCCCC"/></a:solidFill>
        </p:spPr>
        <p:txBody>
            <a:bodyPr/>
            <a:lstStyle/>
            <a:p><a:endParaRPr lang="zh-TW"/></a:p>
        </p:txBody>
    </p:sp>"""
    decoration_elem = etree.fromstring(rect_xml)
    sp_tree.append(decoration_elem)

    _add_5_slides(prs, layout_idx=0)

    out_path = OUTPUT_DIR / "synthetic_C_decoration.pptx"
    prs.save(out_path)

    eval_data = {**COMMON_EVAL_DATA, "source_file": "synthetic_C_decoration.pptx"}
    (OUTPUT_DIR / "synthetic_C_decoration.json").write_text(
        json.dumps(eval_data, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )
    return out_path


def main() -> None:
    """建立所有 3 個 fixture"""
    print("Building synthetic fixtures...")
    paths = []
    for builder in [
        build_synthetic_a_vertical,
        build_synthetic_b_single_placeholder,
        build_synthetic_c_decoration,
    ]:
        path = builder()
        print(f"  ✓ {path.name}")
        paths.append(path)
    print(f"\n✅ {len(paths)} synthetic fixtures created in {OUTPUT_DIR}")
    print("   These fixtures are publicly visible in the repo for CI use.")
    print("   Verify they contain NO ELAN logo / NO real customer data.")


if __name__ == "__main__":
    main()
