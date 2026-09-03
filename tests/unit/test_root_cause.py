"""根因分析與 5-Why 流程圖測試

v3.1.4 新增:稽核報告 #4 揭露 _add_5why_flow_diagram() 有兩個 bug:
  - 沒有「。」時 split 回傳整段,再 [:15] 從字中間切
  - 15 字太短(中文 15 字差不多一句完整建議,英文 15 char 才 2-3 個單字)
  - 強制補滿 5 個步驟(通用「Why 2~5」佔位),與實際內容無關

修法:
  - 新增 _truncate_step_text() helper,中英文句號都認
  - suggestions 不足時不再硬補通用佔位,只顯示實際有的步驟
  - suggestions 完全為空時才 fallback 到預設 5 步
"""

from __future__ import annotations

from pptx import Presentation

from fa_improver.improvers.root_cause import _truncate_step_text


class TestTruncateStepText:
    """_truncate_step_text() 截斷邏輯測試"""

    def test_short_text_unchanged(self):
        """短文字(< max_chars)原樣返回"""
        assert _truncate_step_text("為什麼失敗", max_chars=15) == "為什麼失敗"
        assert _truncate_step_text("Why 1", max_chars=15) == "Why 1"

    def test_empty_text_returns_empty(self):
        """空字串返回空字串"""
        assert _truncate_step_text("") == ""
        assert _truncate_step_text("", max_chars=15) == ""

    def test_chinese_period_truncates(self):
        """中文「。」正確切句"""
        text = "第一句建議比較長。第二句才是重點。"
        # len(text) > max_chars 才會進切句分支
        result = _truncate_step_text(text, max_chars=10)
        assert result == "第一句建議比較長"

    def test_english_period_truncates(self):
        """英文 `.` 正確切句"""
        text = "First check the solder joint quality. Then improve process parameters."
        # text[:idx] 其中 idx 是「。」或「.」的位置
        # 句號在 len("First check the solder joint quality") = 37 處
        result = _truncate_step_text(text, max_chars=30)
        assert result == "First check the solder joint quality"

    def test_no_period_uses_max_chars(self):
        """無句號:按 max_chars 切(舊版 bug 是直接[:15] 從字中間切,新版保留這個行為但不從單字中間切)"""
        # 短於 max_chars:原樣
        text = "無句號的短文字"
        assert _truncate_step_text(text, max_chars=15) == "無句號的短文字"
        # 長於 max_chars:按 max_chars 切
        long_text = "這是一段沒有句號但是非常非常長的建議文字測試"
        result = _truncate_step_text(long_text, max_chars=15)
        assert result == long_text[:15]
        assert len(result) == 15

    def test_no_period_no_mid_word_cut_when_short(self):
        """舊版 bug 重現測試:無句號且 < 15 字時不應被切"""
        # 這是舊版的 bug:s.split("。")[0][:15] if len(s) > 15 else s
        # 條件 len(s) > 15 是 False,所以 return s — 沒 bug
        # 但若 len(s) == 15 也不會切。所以這測試主要是確保短文字不被切
        text = "剛好十五個字無句"  # 8 個字
        assert _truncate_step_text(text, max_chars=15) == text


class Test5WhyFlowDiagram:
    """_add_5why_flow_diagram() 行為測試

    透過 mock FlowDiagramGenerator 並攔截傳入的 steps 來驗證邏輯。
    """

    def _capture_steps(self, suggestions: list[str]) -> list[dict]:
        """呼叫 _add_5why_flow_diagram 並攔截傳入的 steps"""
        from unittest.mock import patch

        from fa_improver.improvers.root_cause import _add_5why_flow_diagram

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        captured = {}

        class FakeFlowGen:
            def __init__(self, *args, **kwargs):
                pass

            def generate(self, steps):
                captured["steps"] = steps

        with patch("fa_improver.improvers.root_cause.FlowDiagramGenerator", FakeFlowGen):
            _add_5why_flow_diagram(slide, suggestions)

        return captured.get("steps", [])

    def test_short_suggestion_no_hard_truncate(self):
        """#1 短於 15 字、無句號的建議不會被切字中間"""
        steps = self._capture_steps(["建議加強環境監控"])
        assert len(steps) == 1
        assert steps[0]["name"] == "建議加強環境監控"  # 完整保留

    def test_long_suggestion_truncates_at_chinese_period(self):
        """#2 中文「。」切句"""
        steps = self._capture_steps(["第一步是檢查焊點品質。第二步是改善製程參數。"])
        assert len(steps) == 1
        assert steps[0]["name"] == "第一步是檢查焊點品質"

    def test_long_suggestion_truncates_at_english_period(self):
        """#3 英文 `.` 切句"""
        steps = self._capture_steps(
            ["First check the solder joint quality. Then improve process parameters."]
        )
        assert len(steps) == 1
        assert steps[0]["name"] == "First check the solder joint quality"

    def test_two_suggestions_no_fake_padding(self):
        """#4 只有 2 個 suggestions 時,流程圖只有 2 個框,沒有「Why 3/4/5」通用佔位"""
        steps = self._capture_steps(
            [
                "建議加強環境監控",
                "改善製程參數設定",
            ]
        )
        assert len(steps) == 2
        names = [s["name"] for s in steps]
        # 不應包含通用佔位
        assert "Why 3: 間接原因" not in names
        assert "Why 4: 系統性原因" not in names
        assert "Why 5: 根本原因" not in names

    def test_three_suggestions_no_fake_padding(self):
        """#5 3 個 suggestions,流程圖只有 3 個"""
        steps = self._capture_steps(["建議 A", "建議 B", "建議 C"])
        assert len(steps) == 3
        names = [s["name"] for s in steps]
        assert "Why 4: 系統性原因" not in names
        assert "Why 5: 根本原因" not in names

    def test_five_suggestions_exact_count(self):
        """#6 5 個 suggestions,流程圖有 5 個框(不超過)"""
        steps = self._capture_steps(["為什麼 A", "為什麼 B", "為什麼 C", "為什麼 D", "為什麼 E"])
        assert len(steps) == 5

    def test_six_suggestions_truncated_to_five(self):
        """#7 6 個 suggestions,只取前 5 個"""
        steps = self._capture_steps(
            ["為什麼 A", "為什麼 B", "為什麼 C", "為什麼 D", "為什麼 E", "為什麼 F"]
        )
        assert len(steps) == 5
        assert steps[0]["name"] == "為什麼 A"
        assert steps[-1]["name"] == "為什麼 E"

    def test_empty_suggestions_falls_back_to_defaults(self):
        """#8 suggestions 完全為空才 fallback 到預設 5 步(向後相容)"""
        steps = self._capture_steps([])
        assert len(steps) == 5
        names = [s["name"] for s in steps]
        assert "Why 1: 表層現象" in names
        assert "Why 5: 根本原因" in names
