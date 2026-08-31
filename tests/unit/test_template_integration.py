"""測試 improvers 與 TemplateLoader 整合

驗證:
- 從 TemplateLoader 載入的樣板標題會用於投影片標題
- placeholder_items 中的 {variable} 會被替換為實際值
- 向後相容(不傳 loader)仍能正常運作
"""

from __future__ import annotations

import json
import tempfile
from pathlib import Path

from fa_improver.domain.evaluation import EvaluationResult
from fa_improver.improvers._template_helper import (
    get_resolved_placeholders,
    resolve_template,
    substitute_placeholders,
)
from fa_improver.improvers.basic_info import add_basic_info_slide
from fa_improver.improvers.prevention import add_prevention_measures_slide
from fa_improver.improvers.root_cause import add_statistical_analysis_slide
from fa_improver.improvers.summary import enhance_summary_section
from fa_improver.parsers.filename_parser import FilenameInfo
from fa_improver.templates.loader import TemplateLoader


def _create_test_pptx():
    """建立測試用 PPTX(用 pptx 預設範本)"""
    from pptx import Presentation

    return Presentation()


def _extract_text_from_slide(slide) -> str:
    """從投影片提取所有文字"""
    texts = []
    for shape in slide.shapes:
        if hasattr(shape, "text_frame"):
            texts.append(shape.text_frame.text)
    return "\n".join(texts)


class TestTemplateHelper:
    """_template_helper 測試"""

    def test_substitute_placeholders(self):
        """替換 {variable} 佔位符"""
        result = substitute_placeholders("FA 編號: {fa_id}", {"fa_id": "FA-001"})
        assert result == "FA 編號: FA-001"

    def test_substitute_placeholders_missing(self):
        """缺少變數時保留原佔位符"""
        result = substitute_placeholders("FA 編號: {fa_id}", {})
        assert result == "FA 編號: {fa_id}"

    def test_substitute_placeholders_empty(self):
        """空字串處理"""
        assert substitute_placeholders("", {"x": "y"}) == ""
        assert substitute_placeholders("no vars here", {"x": "y"}) == "no vars here"

    def test_substitute_multiple_placeholders(self):
        """多個佔位符"""
        text = "{a} and {b}"
        result = substitute_placeholders(text, {"a": "1", "b": "2"})
        assert result == "1 and 2"

    def test_resolve_template_builtin(self):
        """resolve_template 載入內建樣板"""
        template = resolve_template(None, "basic_info")
        assert template.name == "basic_info"
        assert template.title == "FA 基本資訊"

    def test_get_resolved_placeholders(self):
        """從樣板取得 placeholder items 並套用變數"""
        loader = TemplateLoader()
        template = loader.load("basic_info")
        items = get_resolved_placeholders(template, section_index=0, variables={"fa_id": "FA-001"})
        assert len(items) == 7
        assert any("FA-001" in item for item in items)


class TestBasicInfoTemplateIntegration:
    """basic_info 整合 TemplateLoader 測試"""

    def test_uses_template_title(self):
        """使用樣板標題"""
        prs = _create_test_pptx()
        filename_info = FilenameInfo(
            full_stem="FA-001_ACME_X1",
            date_id="001",
            customer="ACME",
            project="X1",
            date="20260831",
        )
        evaluation = EvaluationResult(
            total_score=80.0,
            grade="B",
            dimensions=[],
            summary="",
            strengths=[],
            source_file="test.pptx",
        )
        original_count = len(prs.slides)

        add_basic_info_slide(prs, evaluation, filename_info)

        # 應新增 1 張投影片
        assert len(prs.slides) == original_count + 1
        # 標題應來自樣板
        slide = prs.slides[-1]
        text = _extract_text_from_slide(slide)
        assert "FA 基本資訊" in text

    def test_custom_template_overrides_title(self):
        """自訂樣板覆寫標題"""
        with tempfile.TemporaryDirectory() as tmp:
            tmp_path = Path(tmp)
            custom = {
                "name": "basic_info",
                "title": "【客製】FA 基本資料",
                "layout_name": "2L - Topic",
                "sections": [
                    {
                        "heading": "基本資料",
                        "visual": "bullet_list",
                        "max_bullets": 7,
                        "placeholder_items": ["FA 編號: {fa_id}", "客戶: {customer}"],
                    },
                    {"heading": "優化建議", "max_bullets": 3},
                ],
            }
            (tmp_path / "basic_info.json").write_text(json.dumps(custom), encoding="utf-8")

            loader = TemplateLoader(custom_template_dir=tmp_path)

            prs = _create_test_pptx()
            filename_info = FilenameInfo(
                full_stem="FA-002_X",
                date_id="002",
                customer="ACME",
                project="X",
                date="20260831",
            )
            evaluation = EvaluationResult(
                total_score=80.0,
                grade="B",
                dimensions=[],
                summary="",
                strengths=[],
                source_file="test.pptx",
            )

            add_basic_info_slide(prs, evaluation, filename_info, template_loader=loader)

            text = _extract_text_from_slide(prs.slides[-1])
            assert "【客製】FA 基本資料" in text
            # 變數替換
            assert "FA-002" in text
            assert "ACME" in text

    def test_backward_compatible_without_loader(self):
        """向後相容:不傳 loader 仍正常運作"""
        prs = _create_test_pptx()
        filename_info = FilenameInfo(
            full_stem="FA-003",
            date_id="003",
            customer=None,
            project=None,
            date=None,
        )
        evaluation = EvaluationResult(
            total_score=80.0,
            grade="B",
            dimensions=[],
            summary="",
            strengths=[],
            source_file="test.pptx",
        )

        add_basic_info_slide(prs, evaluation, filename_info)  # 無 loader

        text = _extract_text_from_slide(prs.slides[-1])
        assert "FA 基本資訊" in text
        assert "FA-003" in text


class TestRootCauseTemplateIntegration:
    """root_cause 整合 TemplateLoader 測試"""

    def test_5why_uses_template_title(self):
        """5_why variant 用樣板標題"""
        prs = _create_test_pptx()
        evaluation = EvaluationResult(
            total_score=50.0,
            grade="D",
            dimensions=[],
            summary="",
            strengths=[],
            source_file="test.pptx",
        )
        add_statistical_analysis_slide(
            prs,
            evaluation,
            suggestions=["加強對照組"],
            variant="5_why",
        )
        text = _extract_text_from_slide(prs.slides[-1])
        assert "5-Why 根因推導" in text

    def test_statistical_uses_template_title(self):
        """statistical variant 用樣板標題"""
        prs = _create_test_pptx()
        evaluation = EvaluationResult(
            total_score=50.0,
            grade="D",
            dimensions=[],
            summary="",
            strengths=[],
            source_file="test.pptx",
        )
        add_statistical_analysis_slide(
            prs,
            evaluation,
            suggestions=["加強統計"],
            variant="statistical",
        )
        text = _extract_text_from_slide(prs.slides[-1])
        assert "根因驗證及統計分析" in text

    def test_custom_template_in_root_cause(self):
        """自訂樣板覆寫標題"""
        with tempfile.TemporaryDirectory() as tmp:
            tmp_path = Path(tmp)
            custom = {
                "name": "root_cause_5why",
                "title": "客製 5-Why",
                "layout_name": "2L - Topic",
                "sections": [
                    {
                        "heading": "客製 Heading 1",
                        "visual": "bullet_list",
                        "max_bullets": 4,
                    },
                    {
                        "heading": "客製 Actions",
                        "visual": "checklist",
                        "max_bullets": 3,
                        "placeholder_items": ["客製 action 1", "客製 action 2"],
                    },
                ],
            }
            (tmp_path / "root_cause_5why.json").write_text(json.dumps(custom), encoding="utf-8")
            loader = TemplateLoader(custom_template_dir=tmp_path)

            prs = _create_test_pptx()
            evaluation = EvaluationResult(
                total_score=50.0,
                grade="D",
                dimensions=[],
                summary="",
                strengths=[],
                source_file="test.pptx",
            )
            add_statistical_analysis_slide(
                prs,
                evaluation,
                suggestions=["加強對照組"],
                variant="5_why",
                template_loader=loader,
            )
            text = _extract_text_from_slide(prs.slides[-1])
            assert "客製 5-Why" in text
            assert "客製 Heading 1" in text
            assert "客製 Actions" in text
            assert "客製 action 1" in text


class TestPreventionTemplateIntegration:
    """prevention 整合 TemplateLoader 測試"""

    def test_uses_template_title(self):
        """使用樣板標題"""
        prs = _create_test_pptx()
        evaluation = EvaluationResult(
            total_score=60.0,
            grade="C",
            dimensions=[],
            summary="",
            strengths=[],
            source_file="test.pptx",
        )

        from fa_improver.domain.suggestion import Improvement, Priority

        add_prevention_measures_slide(
            prs,
            evaluation,
            improvements=[
                Improvement(
                    priority=Priority.HIGH,
                    item="對策1",
                    suggestion="建立 IQC SOP",
                )
            ],
        )
        text = _extract_text_from_slide(prs.slides[-1])
        assert "長期預防措施與改善對策" in text
        assert "建立 IQC SOP" in text
        assert "IQC" in text  # 來自樣板 placeholder


class TestSummaryTemplateIntegration:
    """summary 整合 TemplateLoader 測試"""

    def test_uses_template_section_headings(self):
        """使用樣板 section headings"""
        prs = _create_test_pptx()
        # 建立一張 Summary 投影片

        layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        if slide.shapes.title:
            slide.shapes.title.text_frame.text = "Summary"

        from fa_improver.domain.suggestion import Improvement, Priority

        evaluation = EvaluationResult(
            total_score=85.0,
            grade="B",
            dimensions=[],
            summary="報告內容完整,分析詳實",
            strengths=["完整分析 8D 流程", "提供統計驗證"],
            source_file="test.pptx",
        )
        improvements = [Improvement(priority=Priority.HIGH, item="test", suggestion="改進項 A")]

        enhance_summary_section(prs, evaluation, improvements)

        # 注入完成,Summary 投影片應含樣板 section headings
        slide = prs.slides[0]
        text = _extract_text_from_slide(slide)
        assert "Executive Summary" in text
        assert "Key Improvements" in text

    def test_custom_template_changes_headings(self):
        """自訂樣板覆寫 section headings"""
        with tempfile.TemporaryDirectory() as tmp:
            tmp_path = Path(tmp)
            custom = {
                "name": "executive_summary",
                "title": "Summary (客製)",
                "layout_name": "Topic",
                "sections": [
                    {"heading": "原 Summary", "max_bullets": 5},
                    {"heading": "客製優點區", "max_bullets": 5},
                    {"heading": "客製 Executive", "max_bullets": 1},
                    {"heading": "客製 Key Improvements", "max_bullets": 5},
                ],
            }
            (tmp_path / "executive_summary.json").write_text(json.dumps(custom), encoding="utf-8")
            loader = TemplateLoader(custom_template_dir=tmp_path)

            prs = _create_test_pptx()
            from fa_improver.domain.suggestion import Improvement, Priority

            evaluation = EvaluationResult(
                total_score=85.0,
                grade="B",
                dimensions=[],
                summary="",
                strengths=[],
                source_file="test.pptx",
            )
            improvements = [Improvement(priority=Priority.HIGH, item="t", suggestion="s")]

            # 建立 Summary slide
            layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(layout)
            if slide.shapes.title:
                slide.shapes.title.text_frame.text = "Summary"

            enhance_summary_section(prs, evaluation, improvements, template_loader=loader)
            text = _extract_text_from_slide(prs.slides[0])
            assert "客製 Executive" in text
            assert "客製 Key Improvements" in text


class TestOrchestratorTemplateIntegration:
    """Orchestrator 整合 TemplateLoader 測試"""

    def test_orchestrator_accepts_loader(self):
        """Orchestrator 接受 TemplateLoader 參數"""
        from fa_improver.improvers.orchestrator import ImprovementOrchestrator

        evaluation = EvaluationResult(
            total_score=80.0,
            grade="B",
            dimensions=[],
            summary="",
            strengths=[],
            source_file="test.pptx",
        )

        with tempfile.TemporaryDirectory() as tmp:
            loader = TemplateLoader(custom_template_dir=Path(tmp))
            orchestrator = ImprovementOrchestrator(
                evaluation=evaluation,
                input_path=Path("test.pptx"),
                template_loader=loader,
            )
            assert orchestrator.template_loader is loader

    def test_orchestrator_loader_none_by_default(self):
        """Orchestrator 預設 loader 為 None(向後相容)"""
        from fa_improver.improvers.orchestrator import ImprovementOrchestrator

        evaluation = EvaluationResult(
            total_score=80.0,
            grade="B",
            dimensions=[],
            summary="",
            strengths=[],
            source_file="test.pptx",
        )

        orchestrator = ImprovementOrchestrator(evaluation=evaluation, input_path=Path("test.pptx"))
        assert orchestrator.template_loader is None
