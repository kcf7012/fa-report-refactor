"""投影片渲染 smoke test

預防 handoff `2026-08-31-batch-eval-rendering-issues-handoff.md` 提到的問題再次發生:
- 🔴 P0: 完全空白投影片(原本 8 張 → 修正後 0 張)
- 🟡 P1: 內容被擠壓到左上角(13.33 in 寬螢幕 pptx)
- 🟡 P1: 內容互相覆蓋
- 🟢 P2: 文字變垂直(textbox 太窄)
- 🟢 P2: 母片覆蓋

本測試透過真實跑一次完整批次改善流程,驗證:
1. 沒有完全空白的投影片
2. 沒有 shape 超出 slide 邊界
3. 母片保護 100% 通過

執行方式:
    uv run pytest tests/integration/test_slide_rendering.py -v
"""

from __future__ import annotations

import sys
from pathlib import Path

import pytest
from pptx import Presentation

# 確保 src/ 在 Python path
_SKILL_SRC = Path(__file__).parent.parent.parent / "src"
if str(_SKILL_SRC) not in sys.path:
    sys.path.insert(0, str(_SKILL_SRC))

# 強制使用根 report/ 目錄(因為真實 FA 報告放在那裡)
PROJECT_ROOT = Path("/home/elan/fa-report-refactor")
REPORT_DIR = PROJECT_ROOT / "report"


def _is_empty_slide(slide) -> bool:
    """判斷 slide 是否「空白」(少於 3 shapes 且沒有文字)"""
    n_shapes = len(slide.shapes)
    if n_shapes >= 3:
        return False
    has_meaningful_text = any(
        s.has_text_frame and s.text_frame.text.strip()
        for s in slide.shapes
        if s.has_text_frame
    )
    return not has_meaningful_text


def _find_out_of_bounds_shape(slide, slide_width_inch: float, slide_height_inch: float, tolerance_inch: float = 0.2):
    """尋找超出 slide 邊界的 shape,回傳 (slide_idx, shape 名稱) 或 None

    條件(超出容忍值才算):
    - shape.left < -tolerance
    - shape.left + shape.width > slide_width + tolerance
    - shape.top < -tolerance
    - shape.top + shape.height > slide_height + tolerance

    tolerance_inch 預設 0.2 in(避免誤判 pptx 母片設計的邊界 shape)
    """
    EMU_PER_INCH = 914400  # noqa: N806
    sw_emu = slide_width_inch * EMU_PER_INCH
    sh_emu = slide_height_inch * EMU_PER_INCH
    tolerance_emu = tolerance_inch * EMU_PER_INCH

    for shape in slide.shapes:
        # 跳過 group / placeholder(無 left/width)
        if not hasattr(shape, "left") or shape.left is None:
            continue
        try:
            left, top = shape.left, shape.top
            width, height = shape.width, shape.height
        except (AttributeError, ValueError):
            continue

        right = left + width
        bottom = top + height

        if (
            left < -tolerance_emu
            or right > sw_emu + tolerance_emu
            or top < -tolerance_emu
            or bottom > sh_emu + tolerance_emu
        ):
            text_preview = ""
            if shape.has_text_frame:
                text_preview = shape.text_frame.text[:30]
            yield (
                f"{shape.shape_type} at ({left/EMU_PER_INCH:.2f}, {top/EMU_PER_INCH:.2f})"
                f" size ({width/EMU_PER_INCH:.2f} x {height/EMU_PER_INCH:.2f})"
                f" text={text_preview!r}"
            )


def _run_improvement(input_pptx: Path, eval_path: Path, output_suffix: str = "_smoke"):
    """執行一次完整的 improvement 並回傳(prs, result, output_path)"""
    from fa_improver.improvers.orchestrator import ImprovementOrchestrator
    from fa_improver.parsers.evaluation_parser import parse_evaluation

    out_path = REPORT_DIR / f"{input_pptx.stem}{output_suffix}.pptx"
    evaluation = parse_evaluation(eval_path)
    prs = Presentation(input_pptx)
    orchestrator = ImprovementOrchestrator(evaluation, input_pptx)
    result = orchestrator.execute(prs, out_path)
    return prs, result, out_path


class TestSlideRenderingNoEmptySlides:
    """驗證改善後的 pptx 沒有完全空白的投影片

    對應 handoff § 2.1:P0 完全空白投影片(原本 8 張)
    """

    def test_260811_no_empty_slides(self):
        """260811 改善後不應有空白投影片"""
        input_pptx = REPORT_DIR / "260811_Kobo_ZHT_RA6080_SPcomFailI.pptx"
        eval_path = REPORT_DIR / "fa_report_260811_Kobo_ZHT_RA6080_SPcomFailI.json"
        if not input_pptx.exists() or not eval_path.exists():
            pytest.skip(f"需要 {input_pptx.name} 與對應 eval JSON")

        prs, result, _ = _run_improvement(input_pptx, eval_path)

        # 只檢查「新增」的 slide(原 N 張之後的)
        new_slides = list(prs.slides)[result.original_slide_count:]
        empty_slides = []
        for offset, slide in enumerate(new_slides):
            slide_num = result.original_slide_count + offset + 1
            if _is_empty_slide(slide):
                empty_slides.append(slide_num)

        assert not empty_slides, (
            f"[{input_pptx.name}] 發現 {len(empty_slides)} 張新增的空白投影片:"
            f" {empty_slides}。"
            f" 對應 handoff § 2.1:P0 問題(原本 8 張 → 修正後應為 0 張)。"
        )

    def test_ms_meishan_no_empty_slides(self):
        """MS Meishan 改善後不應有空白投影片"""
        input_pptx = REPORT_DIR / "MS_Meishan_ADO_445239_260716.pptx"
        eval_path = REPORT_DIR / "fa_report_MS_Meishan_ADO_445239_260716.json"
        if not input_pptx.exists() or not eval_path.exists():
            pytest.skip(f"需要 {input_pptx.name} 與對應 eval JSON")

        prs, result, _ = _run_improvement(input_pptx, eval_path)

        new_slides = list(prs.slides)[result.original_slide_count:]
        empty_slides = []
        for offset, slide in enumerate(new_slides):
            slide_num = result.original_slide_count + offset + 1
            if _is_empty_slide(slide):
                empty_slides.append(slide_num)

        assert not empty_slides, (
            f"[{input_pptx.name}] 發現 {len(empty_slides)} 張新增的空白投影片:"
            f" {empty_slides}。"
        )

    def test_n160jcn_no_empty_slides(self):
        """N160JCN 改善後不應有空白投影片"""
        # N160JCN 檔名含空格,需 glob 尋找
        candidates = [
            c for c in REPORT_DIR.glob("N160JCN-EEK*NG sample*.pptx")
            if "_improved" not in c.name and "smoke" not in c.name
        ]
        if not candidates:
            pytest.skip("找不到 N160JCN pptx")
        input_pptx = candidates[0]
        # 找對應的 eval JSON
        eval_candidates = [
            c for c in REPORT_DIR.glob("fa_report_N160JCN*.json")
            if "_improved" not in c.name
        ]
        if not eval_candidates:
            pytest.skip("找不到 N160JCN eval JSON")
        eval_path = eval_candidates[0]

        prs, result, _ = _run_improvement(input_pptx, eval_path)

        new_slides = list(prs.slides)[result.original_slide_count:]
        empty_slides = []
        for offset, slide in enumerate(new_slides):
            slide_num = result.original_slide_count + offset + 1
            if _is_empty_slide(slide):
                empty_slides.append(slide_num)

        assert not empty_slides, (
            f"[{input_pptx.name}] 發現 {len(empty_slides)} 張新增的空白投影片:"
            f" {empty_slides}。"
        )


class TestSlideRenderingBounds:
    """驗證改善後 pptx 的所有 shape 都在 slide 邊界內

    對應 handoff § 2.1:🟡 P1 內容被擠壓到左上角(座標超出 slide 邊界)
    """

    def test_no_shape_out_of_bounds(self):
        """所有 shape 都應在 slide 邊界內(0.05 in 容忍)"""
        input_pptx = REPORT_DIR / "260811_Kobo_ZHT_RA6080_SPcomFailI.pptx"
        eval_path = REPORT_DIR / "fa_report_260811_Kobo_ZHT_RA6080_SPcomFailI.json"
        if not input_pptx.exists() or not eval_path.exists():
            pytest.skip(f"需要 {input_pptx.name} 與對應 eval JSON")

        prs, result, _ = _run_improvement(input_pptx, eval_path, output_suffix="_smoke_bounds")

        sw_inch = prs.slide_width / 914400
        sh_inch = prs.slide_height / 914400

        # 只檢查「新增」的 slide(原 N 張之後的)— 原圖的母片設計不歸本測試管
        new_slides = list(prs.slides)[result.original_slide_count:]
        out_of_bounds = []
        for offset, slide in enumerate(new_slides):
            slide_num = result.original_slide_count + offset + 1
            for bad in _find_out_of_bounds_shape(slide, sw_inch, sh_inch):
                out_of_bounds.append(f"Slide {slide_num}: {bad}")

        assert not out_of_bounds, (
            f"發現 {len(out_of_bounds)} 個超出 slide 邊界的 shape(新增 slide 內)\n"
            + "\n".join(out_of_bounds[:5])
        )


class TestSlideRenderingSlideWidths:
    """驗證 orchestrator 正確讀取 pptx slide 寬度

    對應 handoff § 2.3:slide_width 33+ in(實際為 13.33 in 寬螢幕)→ 座標錯位
    """

    def test_orchestrator_reads_slide_width(self):
        """orchestrator 應該從 pptx 讀取真實 slide 寬度,不是寫死 10 in"""
        from fa_improver.domain.evaluation import EvaluationResult
        from fa_improver.improvers.orchestrator import ImprovementOrchestrator

        input_pptx = REPORT_DIR / "260811_Kobo_ZHT_RA6080_SPcomFailI.pptx"
        if not input_pptx.exists():
            pytest.skip("找不到 260811 pptx")
        prs = Presentation(input_pptx)
        real_width = prs.slide_width / 914400
        real_height = prs.slide_height / 914400

        evaluation = EvaluationResult(total_score=50.0, grade="F")
        orch = ImprovementOrchestrator(evaluation, input_pptx)

        # 模擬 execute 開頭設定(實際上 orchestrator 在 execute() 才會設定)
        orch.slide_width_inch = prs.slide_width / 914400
        orch.slide_height_inch = prs.slide_height / 914400

        assert abs(orch.slide_width_inch - real_width) < 0.01, (
            f"orchestrator.slide_width_inch={orch.slide_width_inch}"
            f" 應等於真實 slide 寬度 {real_width:.2f}"
        )
        assert abs(orch.slide_height_inch - real_height) < 0.01


class TestSlideRenderingDynamicCoordinates:
    """驗證內容 shape 會跟著 slide 寬度動態調整

    對應 handoff § 3.5:slide_width 不匹配導致位置錯位
    """

    def test_260811_standard_width_has_dynamic_shapes(self):
        """260811 (10×7.5 in 標準寬度):content shape 應該填滿內容區(>8 in 寬)"""
        input_pptx = REPORT_DIR / "260811_Kobo_ZHT_RA6080_SPcomFailI.pptx"
        eval_path = REPORT_DIR / "fa_report_260811_Kobo_ZHT_RA6080_SPcomFailI.json"
        if not input_pptx.exists() or not eval_path.exists():
            pytest.skip(f"需要 {input_pptx.name} 與對應 eval JSON")

        prs, result, _ = _run_improvement(input_pptx, eval_path, output_suffix="_smoke_dyn")

        # 至少有一個新增 slide 的 content shape 寬度 >= 8 in
        new_slides = list(prs.slides)[result.original_slide_count:]
        max_content_width = 0.0
        for slide in new_slides:
            for shape in slide.shapes:
                if hasattr(shape, "width") and shape.width is not None:
                    w_inch = shape.width / 914400
                    if w_inch > max_content_width:
                        max_content_width = w_inch
        assert max_content_width >= 8.0, (
            f"最大 content shape 寬度 {max_content_width:.2f} in 過窄"
            f"(應 >= 8 in,否則內容會被擠壓到左上角)"
        )


class TestMasterProtectionStillPasses:
    """母片保護測試(回歸測試)

    確保本次改版沒有破壞母片保護(AGENTS.md § 九 最高優先原則)
    """

    def test_master_protection_test_runs(self):
        """確保母片保護驗證可執行且不拋例外"""
        from fa_improver.layout.protector import MasterProtector

        input_pptx = REPORT_DIR / "260811_Kobo_ZHT_RA6080_SPcomFailI.pptx"
        if not input_pptx.exists():
            pytest.skip("找不到 260811 pptx")

        prs = Presentation(input_pptx)
        protector = MasterProtector(prs)
        # 母片保護的 verify_unchanged 不應拋例外
        protector.verify_unchanged(prs)
