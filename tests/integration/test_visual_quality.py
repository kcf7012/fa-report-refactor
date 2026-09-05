"""視覺品質 smoke test — 防止 v3.1.1 的 4 大版面渲染問題再次發生

對應 handoff `2026-09-01-v311-incomplete-rendering-handoff.md`:
- 🔴 Bug 1:enhance_summary_section 疊加覆蓋(原 Summary 內容被新區塊覆蓋)
- 🟡 Bug 2:title placeholder 找錯(「按一下即可新增文字」殘留)
- 🟡 Bug 3:textbox 旋轉 90°(260811 多張)
- 🟡 Bug 4:底部 placeholder 殘留

本測試**不需 pptx 轉圖**,只檢查 pptx XML / shape 屬性就能抓出問題:
- slide.shapes 數量是否爆量(避免 Bug 1)
- textbox.rotation 是否為 0(避免 Bug 3)
- placeholder 預設文字是否被清掉(避免 Bug 2 / Bug 4)

執行方式:
    uv run pytest tests/integration/test_visual_quality.py -v
"""

from __future__ import annotations

import sys
from pathlib import Path

import pytest
from pptx import Presentation

# 確保 src/ 在 Python path
_SKILL_SRC = Path(__file__).parent.parent.parent / "src"
if str(_SKILL_SRC) not in sys.path:
    sys.path.insert(0, str(_SKILL_SRC))

# v3.1.4 修正(稽核 #2):改用動態 fixture resolver,不再硬編專案絕對路徑
from fa_improver.paths import SKILL_ROOT  # noqa: E402
from tests.integration._fixture_resolver import (  # noqa: E402
    SYNTHETIC_FIXTURE_DIR,
    find_project_root,
    get_report_dir,
    resolve_eval_json,
    resolve_input_pptx,
)

PROJECT_ROOT = find_project_root() or SKILL_ROOT
REPORT_DIR = get_report_dir()

RESIDUAL_TITLE_MARKERS = ("按一下", "Click to add", "Click here to add")


def _run_improvement(input_pptx: Path, eval_path: Path, output_suffix: str = "_vq"):
    from fa_improver.improvers.orchestrator import ImprovementOrchestrator
    from fa_improver.parsers.evaluation_parser import parse_evaluation

    out_path = REPORT_DIR / f"{input_pptx.stem}{output_suffix}.pptx"
    evaluation = parse_evaluation(eval_path)
    prs = Presentation(input_pptx)
    orchestrator = ImprovementOrchestrator(evaluation, input_pptx)
    result = orchestrator.execute(prs, out_path)
    return prs, result, out_path


class TestNoSummaryOverlay:
    """Bug 1 修正後測試:enhance_summary_section 不再覆蓋原 Summary

    原本行為:Executive Summary + Key Improvements + ProgressBar 全部疊加在原 Summary
    新行為:新增獨立投影片(在原 Summary 之後)
    """

    def test_summary_section_creates_independent_slides(self):
        """執行 enhance_summary_section 後,新增的 slide 數 = 1(僅 Executive Summary)

        注意:這裡用 mock — 只跑 enhance_summary_section,不跑整個 orchestrator
        """
        input_pptx = resolve_input_pptx("260811_Kobo_ZHT_RA6080_SPcomFailI")
        if not input_pptx:
            pytest.skip("找不到 260811 pptx")

        from fa_improver.domain.evaluation import EvaluationResult
        from fa_improver.domain.suggestion import Improvement, Priority
        from fa_improver.improvers.summary import enhance_summary_section

        prs = Presentation(input_pptx)
        original_count = len(prs.slides)

        evaluation = EvaluationResult(
            total_score=70.0,
            grade="C",
            dimensions=[],
            summary="這是 executive summary 測試",
            strengths=["完整分析 8D"],
            source_file="test.pptx",
        )
        improvements = [Improvement(priority=Priority.HIGH, item="改進項", suggestion="補強 X")]

        enhance_summary_section(prs, evaluation, improvements)

        # 應該新增至少 2 張獨立 slide(Executive + Key Improvements)
        new_count = len(prs.slides) - original_count
        assert new_count >= 2, (
            f"enhance_summary_section 只新增 {new_count} 張 slide,"
            f"應至少 2 張(Executive Summary + Key Improvements)"
        )

    def test_summary_slide_not_overwritten(self):
        """原 Summary slide 的「Summary」標題文字應保留(不被覆蓋)"""
        input_pptx = resolve_input_pptx("260811_Kobo_ZHT_RA6080_SPcomFailI")
        if not input_pptx:
            pytest.skip("找不到 260811 pptx")

        from fa_improver.domain.evaluation import EvaluationResult
        from fa_improver.improvers.summary import enhance_summary_section

        prs = Presentation(input_pptx)

        # 先記錄所有 slide 的文字 hash
        original_texts = []
        for slide in prs.slides:
            text = "\n".join(s.text_frame.text for s in slide.shapes if s.has_text_frame)
            original_texts.append(text)

        evaluation = EvaluationResult(
            total_score=70.0,
            grade="C",
            dimensions=[],
            summary="這是 executive summary 測試",
            strengths=[],
            source_file="test.pptx",
        )

        enhance_summary_section(prs, evaluation, [])

        # 原 N 張的文字應該全部保留(enhance_summary_section 不修改它們)
        for i, original_text in enumerate(original_texts):
            current_text = "\n".join(
                s.text_frame.text for s in prs.slides[i].shapes if s.has_text_frame
            )
            assert current_text == original_text, (
                f"Slide {i + 1} 的文字被改變了!\n"
                f"原:\n{original_text[:200]}\n"
                f"現:\n{current_text[:200]}"
            )


class TestNoTextboxRotation:
    """Bug 3 修正後測試:新增的 textbox 不應被旋轉(除了形狀/表格)"""

    def test_new_textboxes_not_rotated(self):
        """新增 slide 的 textbox 都應 rotation == 0"""
        input_pptx = resolve_input_pptx("260811_Kobo_ZHT_RA6080_SPcomFailI")
        eval_path = resolve_eval_json("260811_Kobo_ZHT_RA6080_SPcomFailI")
        if not input_pptx or not eval_path:
            pytest.skip("需要 260811 pptx 與 eval JSON")

        prs, result, _ = _run_improvement(input_pptx, eval_path, output_suffix="_vqrot")

        # 檢查新增的 slide 上的 textbox
        new_slides = list(prs.slides)[result.original_slide_count :]
        rotated_textboxes = []
        for offset, slide in enumerate(new_slides):
            slide_num = result.original_slide_count + offset + 1
            for shape in slide.shapes:
                if not hasattr(shape, "rotation"):
                    continue
                # 跳過 placeholder(母片設定可能旋轉,但會被 safe_textbox 過濾)
                if shape.is_placeholder:
                    continue
                # 跳過表格
                if shape.has_table:
                    continue
                # shape_type 17 = TEXT_BOX(可旋轉的)
                if getattr(shape, "rotation", 0) != 0:
                    text = shape.text_frame.text[:30] if shape.has_text_frame else ""
                    rotated_textboxes.append(
                        f"Slide {slide_num}: textbox rotation={shape.rotation}, text={text!r}"
                    )

        assert not rotated_textboxes, (
            f"發現 {len(rotated_textboxes)} 個旋轉的 textbox:\n" + "\n".join(rotated_textboxes[:5])
        )


class TestNoResidualPlaceholders:
    """Bug 2 + Bug 4 修正後測試:沒有「按一下即可新增文字」殘留 placeholder"""

    def test_no_residual_placeholders_in_new_slides(self):
        """新增 slide 不應有「按一下」文字的 placeholder"""
        input_pptx = resolve_input_pptx("N160JCN-EEK project 1pcs NG sample analysis report 260810")
        eval_path = resolve_eval_json("N160JCN-EEK project 1pcs NG sample analysis report 260810")
        if not input_pptx or not eval_path:
            pytest.skip("需要 N160JCN pptx 與 eval JSON")

        prs, result, _ = _run_improvement(input_pptx, eval_path, output_suffix="_vqres")

        new_slides = list(prs.slides)[result.original_slide_count :]
        residuals = []
        for offset, slide in enumerate(new_slides):
            slide_num = result.original_slide_count + offset + 1
            for shape in slide.placeholders:
                if not shape.has_text_frame:
                    continue
                text = shape.text_frame.text
                for marker in RESIDUAL_TITLE_MARKERS:
                    if marker in text:
                        residuals.append(f"Slide {slide_num}: placeholder 含「{marker}」")
                        break

        assert not residuals, f"發現 {len(residuals)} 個殘留 placeholder:\n" + "\n".join(
            residuals[:5]
        )


class TestTitlePlaceholderCorrect:
    """Bug 2 修正後測試:title placeholder 的文字應是實際標題,不是 placeholder 預設文字"""

    def test_new_slides_have_meaningful_titles(self):
        """每張新 slide 的 title placeholder 應有實際標題文字"""
        input_pptx = resolve_input_pptx("MS_Meishan_ADO_445239_260716")
        eval_path = resolve_eval_json("MS_Meishan_ADO_445239_260716")
        if not input_pptx or not eval_path:
            pytest.skip("需要 MS pptx 與 eval JSON")

        prs, result, _ = _run_improvement(input_pptx, eval_path, output_suffix="_vqtit")

        new_slides = list(prs.slides)[result.original_slide_count :]
        empty_titles = []
        for offset, slide in enumerate(new_slides):
            slide_num = result.original_slide_count + offset + 1
            # 找 title placeholder(idx=0)
            title_ph = None
            for ph in slide.placeholders:
                if ph.placeholder_format.idx == 0:
                    title_ph = ph
                    break
            if title_ph is None:
                continue  # 沒有 title placeholder 跳過
            text = title_ph.text_frame.text.strip() if title_ph.has_text_frame else ""
            if not text:
                empty_titles.append(f"Slide {slide_num}: title 為空")

        assert not empty_titles, (
            f"發現 {len(empty_titles)} 張新 slide 的 title 是空的:\n" + "\n".join(empty_titles[:5])
        )


# ================================================================
# v3.1.3 新增測試 — Kenny 2026-09-01 回饋的 3 個版面問題
# 參見 docs/handoff/2026-09-01-v313-user-feedback-fixes-handoff.md
# ================================================================


class TestNoTitleDecorationOverlap:
    """v3.1.3 修正後測試:title 不被母片左上裝飾擋住

    Kenny 2026-09-01 反饋:3 份報告的 title 偏左,第一個字被裝飾區擋住。
    修正策略:get_or_create_title 的 fallback safe_textbox 從 left=0.5
    改成 left=1.2 in,避開 x=0.54-0.97 的裝飾區(深藍直條+淺藍色塊)。
    """

    def test_title_textbox_safe_left(self):
        """新增 slide 的 title textbox 的 left >= 1.0 in(避開裝飾區)"""
        input_pptx = resolve_input_pptx("MS_Meishan_ADO_445239_260716")
        eval_path = resolve_eval_json("MS_Meishan_ADO_445239_260716")
        if not input_pptx or not eval_path:
            pytest.skip("需要 MS pptx 與 eval JSON")

        from fa_improver.improvers._safe_shape import TITLE_SAFE_LEFT_INCH

        prs, result, _ = _run_improvement(input_pptx, eval_path, output_suffix="_vqsafe")

        new_slides = list(prs.slides)[result.original_slide_count :]
        # P4:**不再跳過 placeholder**。舊版 `if shape.is_placeholder: continue`
        # 讓這支測試完全看不到「原生 title placeholder 被沿用、left 沒被檢查」
        # 這條路 —— 而那正是最常見的一條路。placeholder 的 `left` 由
        # python-pptx 的 _InheritsDimensions 解析(本層沒設就往 layout / master
        # 取),讀到的已經是有效值,不必自己追繼承鏈;只有兩層都沒設時是 None。
        overlap_titles = []
        for offset, slide in enumerate(new_slides):
            slide_num = result.original_slide_count + offset + 1
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                # 只看 title 類:top 在 1.0 in 內 且 height <= 1.5 in
                # (top >= 1.5 是 body bullet list;height 上限放寬到 1.5,
                #  因為原生 title placeholder 常比 fallback textbox 高,
                #  用 1.0 會把它們整批濾掉 —— 又變成看不見的漏檢)
                if shape.top is None or shape.height is None:
                    continue
                top_inch = shape.top / 914400
                height_inch = shape.height / 914400
                if top_inch > 1.0 or height_inch > 1.5:
                    continue  # 不是 title(top 在 1.0 以上、或高度超過 1.5)
                # 此 shape 是 title
                if shape.left is None:
                    continue  # layout 與 master 都沒設座標,無從判定
                left_inch = shape.left / 914400
                if left_inch < TITLE_SAFE_LEFT_INCH:
                    overlap_titles.append(
                        f"Slide {slide_num}: title left={left_inch:.2f} in,"
                        f"應 >= {TITLE_SAFE_LEFT_INCH} in,"
                        f"text={shape.text_frame.text[:30]!r}"
                    )

        assert not overlap_titles, (
            f"發現 {len(overlap_titles)} 個 title 被裝飾區擋住:\n" + "\n".join(overlap_titles[:5])
        )


class TestTitlePlaceholderNotOverDecoration:
    """P4:原生 title placeholder 也必須避開母片左上裝飾

    背景(計劃書 P4「稽核優先1」):`get_title_placeholder()` 的策略 1
    (`idx == 0`)與策略 2(`TITLE` / `CENTER_TITLE`)一旦命中就直接
    `return ph`,**完全不看 `left` 座標**。`TITLE_SAFE_LEFT_INCH` 只在
    `ph is None` 的 fallback 分支用得到,所以最常見的那條路反而沒有防線。

    本類別**直接指名** `synthetic_C_decoration.pptx`,不透過
    `FIXTURE_FALLBACKS` 的 stem 對應 —— 對應表會隨真實客戶檔在不在位而
    改指到別的檔,那樣這支測試在 Kenny 的機器上跟在 CI 上驗的根本不是
    同一件事,而且降級是靜默的(見 `_fixture_resolver` 的警告)。
    """

    FIXTURE = SYNTHETIC_FIXTURE_DIR / "synthetic_C_decoration.pptx"
    EVAL = SYNTHETIC_FIXTURE_DIR / "synthetic_C_decoration.json"

    def _improve(self, suffix: str):
        if not self.FIXTURE.exists() or not self.EVAL.exists():
            pytest.skip(f"缺少合成 fixture:{self.FIXTURE.name}")
        return _run_improvement(self.FIXTURE, self.EVAL, output_suffix=suffix)

    def test_new_slide_titles_clear_safe_left(self):
        """每一張新投影片的 title(含原生 placeholder)left >= TITLE_SAFE_LEFT_INCH"""
        from fa_improver.improvers._safe_shape import TITLE_SAFE_LEFT_INCH

        prs, result, _ = self._improve("_vqphleft")

        offenders = []
        for offset, slide in enumerate(list(prs.slides)[result.original_slide_count :]):
            slide_num = result.original_slide_count + offset + 1
            for shape in slide.shapes:
                if not shape.has_text_frame or shape.top is None or shape.height is None:
                    continue
                if shape.top / 914400 > 1.0 or shape.height / 914400 > 1.5:
                    continue
                if shape.left is None:
                    continue
                left_inch = shape.left / 914400
                if left_inch < TITLE_SAFE_LEFT_INCH:
                    offenders.append(
                        f"Slide {slide_num}: {shape.name!r} "
                        f"(placeholder={shape.is_placeholder}) left={left_inch:.2f} in "
                        f"< {TITLE_SAFE_LEFT_INCH} in,text={shape.text_frame.text[:30]!r}"
                    )

        assert not offenders, "title 落在安全左界之左:\n" + "\n".join(offenders[:8])

    def test_no_new_shape_overlaps_left_decoration(self):
        """新投影片上沒有任何 shape 與母片「可迴避的左側裝飾」幾何重疊

        比「left >= 常數」更直接:常數調錯時這支會抓到真正的重疊,
        而不是只驗證程式碼與自己的常數一致。

        只採計**可迴避**的裝飾(靠左緣、右緣未過投影片中線)。滿版背景
        與橫幅往右移一寸也閃不開,本來就是 title 要疊在上面的設計元素,
        把它們算進來會讓這支測試變成必紅的假警報。
        """
        prs, result, _ = self._improve("_vqdecov")

        slide_w = prs.slide_width / 914400
        decorations = [
            (
                sh.left / 914400,
                sh.top / 914400,
                (sh.left + sh.width) / 914400,
                (sh.top + sh.height) / 914400,
                sh.name,
            )
            for sh in prs.slide_master.shapes
            if not sh.is_placeholder
            and None not in (sh.left, sh.top, sh.width, sh.height)
            and sh.width / 914400 < slide_w * 0.9
            and sh.left / 914400 <= slide_w * 0.25
            and (sh.left + sh.width) / 914400 < slide_w * 0.5
        ]
        assert decorations, "fixture 應含至少一個可迴避的左側裝飾,否則這支測試什麼都沒驗到"

        overlaps = []
        for offset, slide in enumerate(list(prs.slides)[result.original_slide_count :]):
            slide_num = result.original_slide_count + offset + 1
            for shape in slide.shapes:
                if None in (shape.left, shape.top, shape.width, shape.height):
                    continue
                s_l = shape.left / 914400
                s_t = shape.top / 914400
                s_r = (shape.left + shape.width) / 914400
                s_b = (shape.top + shape.height) / 914400
                for d_l, d_t, d_r, d_b, d_name in decorations:
                    if s_l < d_r and s_r > d_l and s_t < d_b and s_b > d_t:
                        overlaps.append(
                            f"Slide {slide_num}: {shape.name!r} "
                            f"({s_l:.2f},{s_t:.2f})-({s_r:.2f},{s_b:.2f}) "
                            f"壓到母片裝飾 {d_name!r} "
                            f"({d_l:.2f},{d_t:.2f})-({d_r:.2f},{d_b:.2f})"
                        )

        assert not overlaps, "新投影片壓到母片左側裝飾:\n" + "\n".join(overlaps[:8])


class TestBodyClearsLeftDecoration:
    """P4:body 內容也要避開母片左側裝飾(量測後才發現的缺口)

    量測(`scripts/measure_master_decoration.py`)顯示 260811 母片的
    `Picture 14` 佔 x=0.00~1.12、y=1.23~6.38,而舊的 body margin
    `TITLE_SAFE_LEFT_INCH - 0.2` = 1.00 in 比它的右緣還左 0.12 in,
    正文首字會壓在裝飾上。該 layout 沒有 `showMasterSp="0"`,母片
    shape 確實會被渲染。

    這支需要真實客戶檔(合成 fixture 沒有 body 帶的左側裝飾),
    CI 上會 skip —— 這是已知且刻意的。
    """

    def test_body_clears_master_left_decoration(self):
        input_pptx = resolve_input_pptx("260811_Kobo_ZHT_RA6080_SPcomFailI")
        eval_path = resolve_eval_json("260811_Kobo_ZHT_RA6080_SPcomFailI")
        if not input_pptx or not eval_path:
            pytest.skip("需要 260811 pptx 與 eval JSON")
        if input_pptx.parent == SYNTHETIC_FIXTURE_DIR:
            pytest.skip("只有合成 fixture,body 帶沒有左側裝飾可驗")

        prs, result, _ = _run_improvement(input_pptx, eval_path, output_suffix="_vqbodydec")

        slide_w = prs.slide_width / 914400
        decorations = [
            (
                sh.left / 914400,
                sh.top / 914400,
                (sh.left + sh.width) / 914400,
                (sh.top + sh.height) / 914400,
                sh.name,
            )
            for sh in prs.slide_master.shapes
            if not sh.is_placeholder
            and None not in (sh.left, sh.top, sh.width, sh.height)
            and sh.width / 914400 < slide_w * 0.9
            and sh.left / 914400 <= slide_w * 0.25
            and (sh.left + sh.width) / 914400 < slide_w * 0.5
        ]
        if not decorations:
            pytest.skip("此母片沒有可迴避的左側裝飾")

        overlaps = []
        for offset, slide in enumerate(list(prs.slides)[result.original_slide_count :]):
            slide_num = result.original_slide_count + offset + 1
            for shape in slide.shapes:
                if None in (shape.left, shape.top, shape.width, shape.height):
                    continue
                s_l = shape.left / 914400
                s_t = shape.top / 914400
                s_r = (shape.left + shape.width) / 914400
                s_b = (shape.top + shape.height) / 914400
                for d_l, d_t, d_r, d_b, d_name in decorations:
                    if s_l < d_r and s_r > d_l and s_t < d_b and s_b > d_t:
                        overlaps.append(
                            f"Slide {slide_num}: {shape.name!r} "
                            f"({s_l:.2f},{s_t:.2f})-({s_r:.2f},{s_b:.2f}) "
                            f"壓到母片裝飾 {d_name!r} "
                            f"({d_l:.2f},{d_t:.2f})-({d_r:.2f},{d_b:.2f})"
                        )

        assert not overlaps, "新投影片壓到母片左側裝飾:\n" + "\n".join(overlaps[:8])


class TestBodyHasEnoughHeight:
    """v3.1.3 修正後測試:body 有足夠高度容納 heading + bullets

    Kenny 2026-09-01 反饋:MS Page 10/13/14、N160JCN Page 12/15/16 的
    「標題與內容重疊」,因為 layout body placeholder 太矮(0.51 in)。
    修正策略:get_body_placeholder() 當 height < BODY_MIN_HEIGHT_INCH(1.0 in)
    時 fallback 用 safe_textbox 重新建立 body,確保有足夠空間。
    """

    def test_no_overlap_between_title_and_body(self):
        """新 slide 的 body textbox 不應與 title 重疊

        條件:body.top >= title.bottom(不重疊),且 body.height >= 1.0 in
        """
        input_pptx = resolve_input_pptx("N160JCN-EEK project 1pcs NG sample analysis report 260810")
        eval_path = resolve_eval_json("N160JCN-EEK project 1pcs NG sample analysis report 260810")
        if not input_pptx or not eval_path:
            pytest.skip("需要 N160JCN pptx 與 eval JSON")

        from fa_improver.improvers._safe_shape import BODY_MIN_HEIGHT_INCH

        prs, result, _ = _run_improvement(input_pptx, eval_path, output_suffix="_vqbody")

        new_slides = list(prs.slides)[result.original_slide_count :]
        overlap_slides = []
        for offset, slide in enumerate(new_slides):
            slide_num = result.original_slide_count + offset + 1
            # 找 title(textbox 在頂部)
            title_bottom = None
            body_top = None
            body_height = None
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                if shape.top is None or shape.height is None:
                    continue
                top_inch = shape.top / 914400
                height_inch = shape.height / 914400
                if top_inch < 1.5 and height_inch <= 1.2 and title_bottom is None:
                    title_bottom = top_inch + height_inch
                elif top_inch >= 1.5 and height_inch >= 1.0 and body_top is None:
                    body_top = top_inch
                    body_height = height_inch

            # 條件 1:body.top 應 >= title.bottom(不重疊)
            if title_bottom is not None and body_top is not None and body_top < title_bottom:
                overlap_slides.append(
                    f"Slide {slide_num}: title bottom={title_bottom:.2f} > body top={body_top:.2f}"
                )
            # 條件 2:body.height 應 >= BODY_MIN_HEIGHT_INCH
            if body_height is not None and body_height < BODY_MIN_HEIGHT_INCH:
                overlap_slides.append(
                    f"Slide {slide_num}: body height={body_height:.2f} in, 應 >= {BODY_MIN_HEIGHT_INCH} in"
                )

        assert not overlap_slides, f"發現 {len(overlap_slides)} 個 body 問題:\n" + "\n".join(
            overlap_slides[:5]
        )


class TestDimensionChartOptIn:
    """v3.1.3 修正後測試:「6 維度評分分析」slide 預設關閉

    Kenny 2026-09-01 反饋:終端用戶不需要看到「6 維度評分分析」slide
    (這是內部評分指標)。修正策略:enhance_summary_section 預設不產生此 slide,
    需透過 --include-dimension-chart CLI flag 顯式開啟。
    """

    def test_dimension_chart_skipped_by_default(self):
        """include_dimension_chart=False(預設)時,不應新增「6 維度評分分析」slide"""
        input_pptx = resolve_input_pptx("260811_Kobo_ZHT_RA6080_SPcomFailI")
        eval_path = resolve_eval_json("260811_Kobo_ZHT_RA6080_SPcomFailI")
        if not input_pptx or not eval_path:
            pytest.skip("需要 260811 pptx 與 eval JSON")

        prs, result, _ = _run_improvement(input_pptx, eval_path, output_suffix="_vqdim")

        # 檢查所有新 slide,確認沒有「6 維度評分分析」
        new_slides = list(prs.slides)[result.original_slide_count :]
        dim_chart_slides = []
        for offset, slide in enumerate(new_slides):
            slide_num = result.original_slide_count + offset + 1
            for shape in slide.shapes:
                if shape.has_text_frame and "6 維度評分分析" in shape.text_frame.text:
                    dim_chart_slides.append(f"Slide {slide_num}: 發現「6 維度評分分析」slide")
                    break

        assert not dim_chart_slides, (
            f"預設不應出現「6 維度評分分析」slide,但發現 {len(dim_chart_slides)} 個:\n"
            + "\n".join(dim_chart_slides[:3])
        )

    def test_dimension_chart_enabled_with_flag(self):
        """include_dimension_chart=True 時,應新增「6 維度評分分析」slide(opt-in 正常)"""
        input_pptx = resolve_input_pptx("260811_Kobo_ZHT_RA6080_SPcomFailI")
        eval_path = resolve_eval_json("260811_Kobo_ZHT_RA6080_SPcomFailI")
        if not input_pptx or not eval_path:
            pytest.skip("需要 260811 pptx 與 eval JSON")

        from fa_improver.improvers.orchestrator import ImprovementOrchestrator
        from fa_improver.parsers.evaluation_parser import parse_evaluation

        out_path = REPORT_DIR / f"{input_pptx.stem}_vqdim_optin.pptx"
        evaluation = parse_evaluation(eval_path)
        prs = Presentation(input_pptx)
        orchestrator = ImprovementOrchestrator(evaluation, input_pptx, include_dimension_chart=True)
        result = orchestrator.execute(prs, out_path)

        # 確認有「6 維度評分分析」slide
        new_slides = list(prs.slides)[result.original_slide_count :]
        has_dim_chart = False
        for slide in new_slides:
            for shape in slide.shapes:
                if shape.has_text_frame and "6 維度評分分析" in shape.text_frame.text:
                    has_dim_chart = True
                    break
            if has_dim_chart:
                break

        assert has_dim_chart, "include_dimension_chart=True 時,應有「6 維度評分分析」slide"
