"""端對端改善流程測試"""

import pytest
from pptx import Presentation

from fa_improver.improvers.orchestrator import ImprovementOrchestrator
from fa_improver.parsers.evaluation_parser import parse_evaluation
from fa_improver.paths import resolve_report_file


class TestFullWorkflow:
    """端對端改善流程測試"""

    def test_improve_ms_meishan(self, sample_pptx, sample_eval_json, tmp_path):
        """測試 MS Meishan 報告改善"""
        if sample_pptx is None or sample_eval_json is None:
            pytest.skip("範例檔案不存在")

        evaluation = parse_evaluation(sample_eval_json)
        prs = Presentation(sample_pptx)
        original_count = len(prs.slides)
        original_master_xml = prs.slide_masters[0].element.xml

        orchestrator = ImprovementOrchestrator(evaluation, sample_pptx)
        output_path = tmp_path / "improved.pptx"
        result = orchestrator.execute(prs, output_path)

        # 投影片應該增加
        assert result.final_slide_count > original_count

        # 母片保護
        assert result.master_preserved
        assert prs.slide_masters[0].element.xml == original_master_xml

        # 輸出檔案存在
        assert output_path.exists()
        assert output_path.stat().st_size > 0

    def test_improve_n160jcn_from_txt(self, sample_eval_txt, tmp_path):
        """測試從 TXT 評估改善 N160JCN-EEK"""
        # 路徑由 fa_improver.paths 解析(見 P1)。原本用
        # `Path(__file__).parent.parent.parent / "report"` 只到技能包根目錄,
        # 而 N160JCN 檔在外層根倉庫,所以這個測試從來沒有真正跑過。
        n160_pptx = resolve_report_file(
            "N160JCN-EEK project 1pcs NG sample analysis report 260810.pptx"
        )
        if sample_eval_txt is None or n160_pptx is None:
            pytest.skip("範例檔案不存在")

        evaluation = parse_evaluation(sample_eval_txt)
        prs = Presentation(n160_pptx)
        original_count = len(prs.slides)
        original_master_xml = prs.slide_masters[0].element.xml

        orchestrator = ImprovementOrchestrator(evaluation, n160_pptx)
        output_path = tmp_path / "improved_n160jcn.pptx"
        result = orchestrator.execute(prs, output_path)

        # 母片保護
        assert result.master_preserved
        assert prs.slide_masters[0].element.xml == original_master_xml

        # 因為基本資訊 85 >= 80,不應新增基本資訊
        # 但根因分析 40 < 80,應新增
        # 改善對策 40 < 85,應新增
        assert result.final_slide_count > original_count
