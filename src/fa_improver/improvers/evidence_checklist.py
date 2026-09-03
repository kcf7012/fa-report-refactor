"""數據與證據支持 Improver

針對「數據與證據支持」維度(權重 20%)。
改善內容:
- 對照組 vs 異常品 數據對照表
- 圖片品質檢查清單
- 數據追溯性指引
"""

from __future__ import annotations

from pptx import Presentation

from ..domain.evaluation import DimensionScore, EvaluationResult
from ..layout.selector import find_content_layout
from ..templates.loader import TemplateLoader
from ..visuals import (
    ELAN_BLUE,
    ELAN_GREEN,
    ELAN_ORANGE,
    ELAN_RED,
    ChecklistGenerator,
    ComparisonTableGenerator,
)
from ._safe_shape import TITLE_SAFE_LEFT_INCH
from ._template_helper import resolve_template


def add_evidence_checklist_slide(
    prs: Presentation,
    evaluation: EvaluationResult,
    dimension: DimensionScore | None = None,
    template_loader: TemplateLoader | None = None,
    template_name: str = "evidence_checklist",
    slide_bounds: dict | None = None,
) -> None:
    """新增數據與證據支持投影片

    Args:
        prs: 簡報物件
        evaluation: 評估結果
        dimension: 該維度的評分(可選)
        template_loader: 樣板載入器(可選)
        template_name: 樣板名稱(預設 'evidence_checklist')
        slide_bounds: slide 尺寸(英寸),動態適應
    """
    from ._logging import log_action

    with log_action("add_evidence_checklist_slide"):
        layout = find_content_layout(prs)
        slide = prs.slides.add_slide(layout)

        # === 動態座標 ===
        sw = slide_bounds["width_inch"] if slide_bounds else 10.0
        # v3.1.4:與 _safe_shape safe_textbox fallback 對齊(margin=1.0,確保小 slide 不會太擠)
        margin = TITLE_SAFE_LEFT_INCH - 0.2
        if margin < 0.5:
            margin = 0.5
        content_w = sw - 2 * margin

        # 標題 — 優先從樣板讀取
        title = _get_or_create_title(slide, slide_bounds)
        if template_loader is not None:
            try:
                template = resolve_template(template_loader, template_name)
                title.text_frame.text = template.title
            except KeyError:
                title.text_frame.text = "數據與證據支持"
        else:
            title.text_frame.text = "數據與證據支持"

        # 1. 對照組 vs 異常品 數據對照表
        _add_comparison_data_table(slide, content_w)

        # 2. 圖片品質與數據追溯性檢查清單
        _add_evidence_checklist(slide, content_w)


def _add_comparison_data_table(slide, content_w: float = 9.0) -> None:
    """加入對照組 vs 異常品數據對照表"""
    gen = ComparisonTableGenerator(
        slide,
        left=TITLE_SAFE_LEFT_INCH - 0.2,  # v3.1.4:與 _safe_shape 對齊
        top=1.4,
        width=content_w,
        height=2.2,
    )
    gen.generate(
        {
            "headers": ["測試項目", "DVT 正常品", "PVT 異常品", "Spec 範圍", "判定"],
            "rows": [
                ["ESD HBM (±2kV)", "通過", "?", "±2kV", "需補測"],
                ["I/O 對地阻抗", "16MΩ", "5.7KΩ", ">1MΩ", "FAIL"],
                ["VH/VOUT 電壓", "正常", "0V", "0.4-0.6V", "FAIL"],
                ["二極體特性", "0.43V", "0V", "0.4-0.6V", "FAIL"],
                ["FW 讀取", "正確", "正確", "100%", "PASS"],
                ["外觀檢查", "正常", "正常", "無損傷", "PASS"],
            ],
        }
    )


def _add_evidence_checklist(slide, content_w: float = 9.0) -> None:
    """加入圖片品質與數據追溯性檢查清單"""
    gen = ChecklistGenerator(
        slide,
        left=TITLE_SAFE_LEFT_INCH - 0.2,  # v3.1.4:與 _safe_shape 對齊
        top=3.9,
        width=content_w,
        height=2.8,
    )
    gen.generate(
        [
            {
                "text": "圖片解析度 ≥ 1000x(建議 ≥ 2000x 以利 IC Marking 辨識)",
                "checked": False,
                "color": ELAN_RED,
            },
            {
                "text": "對焦清晰,IC 標籤、Marking 可清楚辨識",
                "checked": False,
                "color": ELAN_RED,
            },
            {
                "text": "含比例尺(scale bar),方便評估損傷範圍",
                "checked": False,
                "color": ELAN_ORANGE,
            },
            {
                "text": "異常點明確標註(箭頭、方框、文字說明)",
                "checked": False,
                "color": ELAN_ORANGE,
            },
            {
                "text": "對照組(Golden Sample)與異常品並列比較",
                "checked": False,
                "color": ELAN_BLUE,
            },
            {
                "text": "每個量化數據有來源追溯(儀器型號、測試條件、日期)",
                "checked": False,
                "color": ELAN_BLUE,
            },
            {
                "text": "統計顯著性(p-value、CI)支援結論(非僅描述)",
                "checked": False,
                "color": ELAN_GREEN,
            },
        ]
    )


def _get_or_create_title(slide, slide_bounds: dict | None = None):
    """取得真實的 title placeholder(Bug 2 + Bug 3 修正)"""
    from ._safe_shape import get_or_create_title

    return get_or_create_title(slide, slide_bounds)
