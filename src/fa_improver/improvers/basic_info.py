"""新增 FA 基本資訊投影片

從 TemplateLoader 載入 'basic_info' 樣板取得標題與 placeholder items。
向後相容:若不傳 loader,使用預設載入器。

視覺元素:使用 ChecklistGenerator 呈現基本資料的檢查狀態。
"""

from __future__ import annotations

from typing import TYPE_CHECKING

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from ..domain.evaluation import EvaluationResult
from ..layout.selector import find_content_layout
from ..parsers.filename_parser import FilenameInfo
from ..templates.loader import TemplateLoader
from ..visuals import ELAN_BLUE, ChecklistGenerator
from ._safe_shape import BODY_SAFE_LEFT_INCH
from ._template_helper import get_resolved_placeholders, resolve_template

if TYPE_CHECKING:
    pass


def add_basic_info_slide(
    prs: Presentation,
    evaluation: EvaluationResult,
    filename_info: FilenameInfo,
    template_loader: TemplateLoader | None = None,
    template_name: str = "basic_info",
    slide_bounds: dict | None = None,
) -> None:
    """新增 FA 基本資訊投影片

    Args:
        prs: 簡報物件
        evaluation: 評估結果
        filename_info: 檔名解析結果
        template_loader: 樣板載入器(可選,預設使用內建樣板)
        template_name: 樣板名稱(預設 'basic_info')
        slide_bounds: slide 尺寸(英寸),若提供則動態適應;否則用 10x7.5 預設
    """
    from ._logging import log_action

    with log_action("add_basic_info_slide"):
        _add_basic_info_slide_impl(
            prs, evaluation, filename_info, template_loader, template_name, slide_bounds
        )


def _add_basic_info_slide_impl(
    prs: Presentation,
    evaluation: EvaluationResult,
    filename_info: FilenameInfo,
    template_loader: TemplateLoader | None,
    template_name: str,
    slide_bounds: dict | None,
) -> None:
    layout = find_content_layout(prs)
    slide = prs.slides.add_slide(layout)

    # === 動態座標(v3.1.1 修正)===
    sw = slide_bounds["width_inch"] if slide_bounds else 10.0
    # v3.1.4:與 _safe_shape safe_textbox fallback 對齊(margin=1.0,確保小 slide 不會太擠)
    margin = BODY_SAFE_LEFT_INCH
    if margin < 0.5:
        margin = 0.5
    content_w = sw - 2 * margin

    # 載入樣板
    template = resolve_template(template_loader, template_name)

    # 標題(從樣板)
    title_shape = _get_or_create_title(slide, slide_bounds)
    title_shape.text_frame.text = template.title

    # 從樣板取得 placeholder items 並套用變數替換
    variables = {
        "fa_id": filename_info.to_fa_id(),
        "customer": filename_info.customer or "N/A",
        "project": filename_info.project or "N/A",
        "date": filename_info.date or "N/A",
    }
    placeholders = get_resolved_placeholders(template, section_index=0, variables=variables)

    # 使用 ChecklistGenerator 呈現基本資料的檢查狀態
    checklist_items = [
        {"text": item, "checked": False, "color": ELAN_BLUE} for item in placeholders
    ]
    if checklist_items:
        checklist_gen = ChecklistGenerator(
            slide,
            left=margin,
            top=1.4,
            width=content_w,
            height=4.0,
        )
        checklist_gen.generate(checklist_items)

    # 優化建議項目(從 comment 抽取,屬於 section index 1)
    if evaluation.dimensions:
        dim = next((d for d in evaluation.dimensions if d.name.value == "基本資訊完整性"), None)
        if dim and dim.comment:
            # section 1 的標題(從樣板)
            section1 = template.sections[1] if len(template.sections) > 1 else None
            heading = section1.heading if section1 else "優化建議項目"

            # 用 textbox 顯示優化建議(在 checklist 下方)
            # 使用 safe_textbox 避免 Bug 3(textbox 旋轉)
            from ._safe_shape import safe_textbox

            body = safe_textbox(
                slide,
                left=margin,
                top=5.5,
                width=content_w,
                height=1.5,
            )
            tf = body.text_frame

            p = tf.paragraphs[0]
            p.text = f"[{heading}]"
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 0, 0)
            p.font.size = Pt(14)

            sub = tf.add_paragraph()
            sub.text = f"• {dim.comment}"
            sub.font.size = Pt(12)


def _get_or_create_title(slide, slide_bounds: dict | None = None):
    """取得真實的 title placeholder(Bug 2 + Bug 3 修正,v3.1.3 統一改用 helper)

    v3.1.3 修正:原本 hard-coded margin=0.5,會被母片左上裝飾(0.54-0.97 in)擋住。
    改為統一呼叫 `_safe_shape.get_or_create_title`,fallback safe_textbox 的 left=1.2 in。
    """
    from ._safe_shape import get_or_create_title

    return get_or_create_title(slide, slide_bounds)


def _get_or_create_body(slide, slide_bounds: dict | None = None):
    """取得 body placeholder,優先選擇 rotation == 0 的(Bug 3 修正)

    260811 的某些 layout body placeholder 被預設為旋轉 90°,
    若直接寫入文字,整頁會變直式。

    v3.1.3 修正:統一改用 `_safe_shape.get_or_create_body`,
    fallback 條件改為:layout placeholder 高度 < BODY_MIN_HEIGHT_INCH(1.0) 時,
    才會 fallback 用 safe_textbox,避免內容與 title 區重疊。
    """
    from ._safe_shape import get_or_create_body

    return get_or_create_body(slide, slide_bounds)


# Helper aliases(縮短程式碼)
_Inches = Inches


def _PT(size: int):  # noqa: N802
    from pptx.util import Pt

    return Pt(size)


def _COLOR(r: int, g: int, b: int):  # noqa: N802
    from pptx.dml.color import RGBColor

    return RGBColor(r, g, b)
