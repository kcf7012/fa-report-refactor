"""Summary 強化 — 展開為多張獨立投影片

從 TemplateLoader 載入 'executive_summary' 樣板取得標題與 sections 標題。
向後相容:若不傳 loader,使用預設載入器。

修正 Bug 1(v3.1.1 殘留問題):
原本 enhance_summary_section 把 Executive Summary + Key Improvements + ProgressBar
全部疊加在原 Summary 投影片上,導致互相覆蓋。
改為新增獨立投影片(在原 Summary 之後新增 3 張 slide)。
"""

from __future__ import annotations

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt

from ..domain.evaluation import EvaluationResult
from ..domain.suggestion import Improvement
from ..templates.loader import TemplateLoader
from ..visuals import (
    ELAN_BLUE,
    ELAN_GREEN,
    ELAN_ORANGE,
    ELAN_RED,
    ProgressBarGenerator,
)
from ._safe_shape import TITLE_SAFE_LEFT_INCH, safe_textbox
from ._template_helper import resolve_template


def enhance_summary_section(
    prs: Presentation,
    evaluation: EvaluationResult,
    improvements: list[Improvement],
    template_loader: TemplateLoader | None = None,
    template_name: str = "executive_summary",
    slide_bounds: dict | None = None,
    *,
    include_dimension_chart: bool = False,
) -> None:
    """強化 Summary 區塊 — 改為新增獨立投影片(不再疊加)

    Bug 1 修正:原實作把 3 個區塊(Executive Summary + Key Improvements +
    ProgressBar)全部疊加在原 Summary 投影片上,造成互相覆蓋。

    新實作:在原 Summary 之後新增 2-3 張獨立投影片:
    1. Executive Summary slide
    2. Key Improvements Required slide
    3. (可選)6 維度評分進度條 slide — 預設關閉(v3.1.3)

    原 Summary 投影片不被修改。

    Args:
        prs: 簡報物件
        evaluation: 評估結果
        improvements: 改進建議
        template_loader: 樣板載入器(可選)
        template_name: 樣板名稱(預設 'executive_summary')
        slide_bounds: slide 尺寸(英寸),動態適應
        include_dimension_chart: 是否包含「6 維度評分分析」slide。
            預設 False(v3.1.3 起的用戶預設 — Kenny 反饋此 slide
            對終端用戶無實質幫助,屬內部評分指標)。
    """
    from ._logging import log_action

    with log_action("enhance_summary_section"):
        # 載入樣板(讀取標題結構)
        try:
            template = resolve_template(template_loader, template_name)
        except KeyError:
            template = None  # fallback

        # Bug 1 修正:新增獨立投影片,而非疊加原 Summary
        # 注意:順序很重要 — Executive Summary → Key Improvements → Dimension Progress
        # 這些 slide 會被加在原 Summary 之後的位置

        # 若有 evaluation.summary(非空)才新增 Executive Summary slide
        # 原本是 is None 才 fallback,現改為空字串也 fallback
        if evaluation.summary is not None:
            _new_executive_summary_slide(prs, evaluation, template, slide_bounds)

        if improvements:
            _new_key_improvements_slide(prs, improvements, evaluation, template, slide_bounds)

        # v3.1.3:6 維度評分分析改為 opt-in,預設關閉
        if include_dimension_chart and evaluation.dimensions:
            _new_dimension_progress_slide(prs, evaluation, slide_bounds)


def _new_executive_summary_slide(
    prs: Presentation,
    evaluation: EvaluationResult,
    template=None,
    slide_bounds: dict | None = None,
) -> None:
    """新增 Executive Summary 獨立 slide

    Args:
        prs: 簡報物件
        evaluation: 評估結果
        template: 樣板(可選)
        slide_bounds: slide 尺寸
    """
    from ..layout.selector import find_content_layout

    # 從樣板讀取 section heading
    heading = "Executive Summary"
    if template and len(template.sections) >= 3:
        heading = template.sections[2].heading or heading

    layout = find_content_layout(prs)
    slide = prs.slides.add_slide(layout)

    sw = slide_bounds["width_inch"] if slide_bounds else 10.0
    # v3.1.4:與 _safe_shape safe_textbox fallback 對齊(margin=1.0,確保小 slide 不會太擠)
    margin = TITLE_SAFE_LEFT_INCH - 0.2
    if margin < 0.5:
        margin = 0.5
    content_w = sw - 2 * margin

    # Title(Bug 2 + Bug 3 修正:用 get_title_placeholder + safe_textbox)
    from ._safe_shape import clean_unused_placeholders, get_or_create_title

    title = get_or_create_title(slide, slide_bounds)
    title.text_frame.text = heading

    # Body 內容
    body = safe_textbox(
        slide,
        left=margin,
        top=1.5,
        width=content_w,
        height=sh - 2.0 if (sh := slide_bounds["height_inch"] if slide_bounds else 7.5) else 5.5,
    )
    tf = body.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = evaluation.summary or "報告分析詳實,建議補充統計數據以強化結論。"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # Bug 4 修正:清掉殘留 placeholder
    clean_unused_placeholders(slide)


def _new_key_improvements_slide(
    prs: Presentation,
    improvements: list[Improvement],
    evaluation: EvaluationResult,
    template=None,
    slide_bounds: dict | None = None,
) -> None:
    """新增 Key Improvements Required 獨立 slide"""
    from ..layout.selector import find_content_layout

    heading = "Key Improvements Required"
    if template and len(template.sections) >= 4:
        heading = template.sections[3].heading or heading

    layout = find_content_layout(prs)
    slide = prs.slides.add_slide(layout)

    sw = slide_bounds["width_inch"] if slide_bounds else 10.0
    # v3.1.4:與 _safe_shape safe_textbox fallback 對齊(margin=1.0,確保小 slide 不會太擠)
    margin = TITLE_SAFE_LEFT_INCH - 0.2
    if margin < 0.5:
        margin = 0.5
    content_w = sw - 2 * margin
    sh = slide_bounds["height_inch"] if slide_bounds else 7.5

    from ._safe_shape import clean_unused_placeholders, get_or_create_title

    title = get_or_create_title(slide, slide_bounds)
    title.text_frame.text = heading

    # Body — 條列改善建議
    body = safe_textbox(
        slide,
        left=margin,
        top=1.5,
        width=content_w,
        height=sh - 2.0,
    )
    tf = body.text_frame
    tf.word_wrap = True

    # 第一段是 heading(已用 title,這裡改為空白第一行)
    for idx, imp in enumerate(improvements[:6]):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"• {imp.suggestion}"
        p.font.size = Pt(13)
        # Priority 是 str enum(HIGH/MEDIUM/LOW),直接用 name 比較
        if hasattr(imp, "priority") and imp.priority.name == "HIGH":
            p.font.color.rgb = RGBColor(255, 0, 0)  # HIGH → 紅色
        else:
            p.font.color.rgb = RGBColor(0, 112, 192)  # 其他 → 藍色

    # Bug 4 修正:清掉殘留 placeholder
    clean_unused_placeholders(slide)


def _new_dimension_progress_slide(
    prs: Presentation,
    evaluation: EvaluationResult,
    slide_bounds: dict | None = None,
) -> None:
    """新增 6 維度評分進度條獨立 slide"""
    from ..layout.selector import find_content_layout

    if not evaluation.dimensions:
        return

    layout = find_content_layout(prs)
    slide = prs.slides.add_slide(layout)

    sw = slide_bounds["width_inch"] if slide_bounds else 10.0
    # v3.1.4:與 _safe_shape safe_textbox fallback 對齊(margin=1.0,確保小 slide 不會太擠)
    margin = TITLE_SAFE_LEFT_INCH - 0.2
    if margin < 0.5:
        margin = 0.5
    content_w = sw - 2 * margin

    from ._safe_shape import clean_unused_placeholders, get_or_create_title

    title = get_or_create_title(slide, slide_bounds)
    title.text_frame.text = "6 維度評分分析"

    # 6 維度進度條
    def color_for_score(score: float):
        if score >= 85:
            return ELAN_GREEN
        if score >= 70:
            return ELAN_BLUE
        if score >= 50:
            return ELAN_ORANGE
        return ELAN_RED

    items = [
        {
            "label": dim_score.name.value,
            "value": dim_score.score,
            "max_value": 100,
            "color": color_for_score(dim_score.score),
        }
        for dim_score in evaluation.dimensions
    ]

    gen = ProgressBarGenerator(
        slide,
        left=1.0,
        top=1.8,
        width=content_w,
        height=4.5,
    )
    gen.generate(items)

    # Bug 4 修正:清掉殘留 placeholder
    clean_unused_placeholders(slide)
