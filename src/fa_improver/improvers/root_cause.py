"""新增根因分析相關投影片

從 TemplateLoader 載入 'root_cause_5why' 或 'root_cause_statistical' 樣板取得標題。
向後相容:若不傳 loader,使用預設載入器。

視覺元素:
- 5-Why variant 使用 FlowDiagramGenerator 呈現推導流程
- statistical variant 使用文字建議
"""

from __future__ import annotations

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from ..domain.evaluation import EvaluationResult
from ..layout.selector import find_content_layout
from ..templates.loader import TemplateLoader
from ..visuals import FlowDiagramGenerator
from ._template_helper import get_resolved_placeholders, resolve_template


def add_statistical_analysis_slide(
    prs: Presentation,
    evaluation: EvaluationResult,
    suggestions: list[str],
    variant: str = "statistical",
    template_loader: TemplateLoader | None = None,
    slide_bounds: dict | None = None,
) -> None:
    """新增根因分析投影片

    Args:
        prs: 簡報物件
        evaluation: 評估結果
        suggestions: 建議清單
        variant:
            - "5_why": 5-Why 推導流程
            - "statistical": 統計驗證方法
        template_loader: 樣板載入器(可選)
        slide_bounds: slide 尺寸(英寸),動態適應
    """
    from ._logging import log_action

    with log_action(f"add_statistical_analysis_slide:variant={variant}"):
        # === 動態座標 ===
        sw = slide_bounds["width_inch"] if slide_bounds else 10.0
        margin = 0.5
        content_w = sw - 2 * margin

        # 載入對應樣板
        template_name = "root_cause_5why" if variant == "5_why" else "root_cause_statistical"
        template = resolve_template(template_loader, template_name)

        layout = find_content_layout(prs)
        slide = prs.slides.add_slide(layout)

        # 標題(從樣板)
        title = _get_or_create_title(slide, slide_bounds)
        title.text_frame.text = template.title

        # 內容
        if not suggestions:
            suggestions = ["建議加強對照組設定與數據統計驗證以支撐根因發現。"]

        # 5_Why variant:加入 FlowDiagramGenerator 視覺化推導流程
        if variant == "5_why":
            _add_5why_flow_diagram(slide, suggestions, content_w)

        body = _get_or_create_body(slide, slide_bounds)
        tf = body.text_frame
        tf.clear()

    # section 0 的標題
    section0 = template.sections[0] if template.sections else None
    heading0 = section0.heading if section0 else "針對問題點之深度分析建議"

    p = tf.paragraphs[0]
    p.text = f"{heading0}:"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(0, 112, 192)

    # 限制在 max_bullets
    max_bullets = section0.max_bullets if section0 else 4
    for sug in suggestions[:max_bullets]:
        p = tf.add_paragraph()
        p.text = sug
        p.font.size = Pt(14)

    # section 1 的標題(若存在)
    section1 = template.sections[1] if len(template.sections) > 1 else None
    if section1:
        p = tf.add_paragraph()
        p.text = f"\n[{section1.heading}]"
        p.font.bold = True

        # 從樣板讀取 placeholder_items
        placeholders = get_resolved_placeholders(template, section_index=1)
        if placeholders:
            for item in placeholders:
                p = tf.add_paragraph()
                p.text = f"• {item}"
                p.font.size = Pt(12)
        else:
            # fallback:統計驗證預設 actions
            for action in [
                "設定 DVT 正常品 vs PVT 異常品之對照組",
                "使用獨立樣本 t 檢定驗證參數顯著性 (p < 0.05)",
                "確保統計證據支持最終提到的根本原因",
            ]:
                p = tf.add_paragraph()
                p.text = f"• {action}"
                p.font.size = Pt(12)


def _truncate_step_text(text: str, max_chars: int = 15) -> str:
    """簡化 step 文字:按句號(中英文)切,保留第一句。

    v3.1.4 修正(稽核 #4):
    - 舊: `s.split("。")[0][:15] if len(s) > 15 else s` 在沒有「。」時
      會返回整段後再 `[:15]`,從字中間切。
    - 新: 同時認中英文句號 `。` 與 `.`;若無句號,保留整段(不超過 max_chars)。
    """
    if not text:
        return ""
    if len(text) <= max_chars:
        return text
    # 按中英文句號切
    for sep in ("。", "."):
        idx = text.find(sep)
        if idx > 0:
            return text[:idx]
    # 無句號:按 max_chars 切
    return text[:max_chars]


def _add_5why_flow_diagram(
    slide, content_w_or_suggestions=None, content_w: float | None = None
) -> None:
    """加入 5-Why 推導流程圖

    從 suggestions 建立流程圖;若 suggestions 不足 5 個,則只顯示實際有的數量,
    而非硬補通用「Why 2~5」佔位符。

    v3.1.4 修正(稽核 #4):
    - 舊:當 suggestions 不足時,把通用 `Why 2: 直接原因` 等佔位補滿 5 個,
      造成流程圖中後幾個框與前幾個毫無關聯。
    - 新:優先使用 suggestions(中英句號都認,避免從字中間切);
      若 suggestions 為空才 fallback 到預設 5 步。

    向後相容:可接受 (slide, suggestions) 或 (slide, content_w, suggestions) 三種呼叫方式。
    """
    # 解析參數
    if content_w is None:
        # 舊簽名: (slide, suggestions)
        suggestions = content_w_or_suggestions or []
        cw = 9.0
    else:
        # 新簽名: (slide, suggestions, content_w)
        suggestions = content_w_or_suggestions or []
        cw = content_w

    # 構建 steps:優先使用 suggestions,不足才 fallback
    if suggestions:
        # 截斷每個 suggestion 並建立 step dict
        steps = [
            {"name": _truncate_step_text(s, max_chars=15), "status": "active"}
            for s in suggestions[:5]  # 最多 5 個
        ]
    else:
        # suggestions 為空:fallback 到預設 5 個步驟
        default_steps = [
            "Why 1: 表層現象",
            "Why 2: 直接原因",
            "Why 3: 間接原因",
            "Why 4: 系統性原因",
            "Why 5: 根本原因",
        ]
        steps = [{"name": s, "status": "active"} for s in default_steps]

    flow_gen = FlowDiagramGenerator(
        slide,
        left=0.5,
        top=4.5,
        width=cw,
        height=2.5,
    )
    flow_gen.generate(steps)


def add_5why_slide(
    prs: Presentation,
    evaluation: EvaluationResult,
    suggestions: list[str],
    variant: str = "5_why",
    template_loader: TemplateLoader | None = None,
    slide_bounds: dict | None = None,
) -> None:
    """5-Why 變體的快捷函式

    為 orchestrator 提供一個統一的進入點;目前 variant 支援:
    - "5_why": 同 add_statistical_analysis_slide(variant="5_why")
    - "control_group": 控制組/對照組分析
    - "evidence": 證據型根因分析

    Args:
        prs: 簡報物件
        evaluation: 評估結果
        suggestions: 建議清單
        variant: "5_why" | "control_group" | "evidence"
        template_loader: 樣板載入器
        slide_bounds: slide 尺寸(英寸)
    """
    from ._logging import get_logger, log_action

    logger = get_logger()
    with log_action(f"add_5why_slide:variant={variant}"):
        if variant == "5_why":
            add_statistical_analysis_slide(
                prs,
                evaluation=evaluation,
                suggestions=suggestions,
                variant="5_why",
                template_loader=template_loader,
                slide_bounds=slide_bounds,
            )
            return

        # 其他 variant:建立獨立的副投影片
        from ..layout.selector import find_content_layout
        from ..visuals import ComparisonTableGenerator

        sw = slide_bounds["width_inch"] if slide_bounds else 10.0
        margin = 0.5
        content_w = sw - 2 * margin

        layout = find_content_layout(prs)
        slide = prs.slides.add_slide(layout)

        title = _get_or_create_title(slide)
        if variant == "control_group":
            title.text_frame.text = "控制組 / 對照組分析"
            headers = ["對照組設定", "目的", "驗證邏輯"]
            rows = [
                ["Golden Sample", "建立失效基準", "與異常品並列比對"],
                ["DVT 正常品", "製程能力驗證", "統計差異分析 (t-test)"],
                ["同批 PVT", "排除批次性失效", "批次間變異數分析"],
                ["跨批抽樣", "長期趨勢監測", "管制圖 (Control Chart)"],
            ]
        elif variant == "evidence":
            title.text_frame.text = "證據型根因分析"
            headers = ["證據類型", "取得方式", "說服力"]
            rows = [
                ["電性數據", "I/V 曲線、阻抗", "高"],
                ["結構影像", "SEM、X-ray", "高"],
                ["成分分析", "EDX、XRD", "中"],
                ["失效重現", "HTSL、HAST", "高"],
                ["統計分析", "p-value、CI", "中"],
            ]
        else:
            logger.warning("未知的 variant: %s", variant)
            return

        gen = ComparisonTableGenerator(
            slide,
            left=margin,
            top=1.4,
            width=content_w,
            height=5.0,
        )
        gen.generate({"headers": headers, "rows": rows})

        # 底部說明
        from pptx.dml.color import RGBColor
        from pptx.util import Pt as _Pt

        note_box = slide.shapes.add_textbox(
            Inches(margin), Inches(6.6), Inches(content_w), Inches(0.5)
        )
        tf = note_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "✓ 此頁為 v3.1.1 新增,提供控制組/對照組/證據型根因分析的標準範本"
        p.font.color.rgb = RGBColor(0x80, 0x80, 0x80)
        p.font.size = _Pt(10)


def _get_or_create_title(slide, slide_bounds: dict | None = None):
    """取得真實的 title placeholder(Bug 2 + Bug 3 修正)"""
    from ._safe_shape import get_or_create_title

    return get_or_create_title(slide, slide_bounds)


def _get_or_create_body(slide, slide_bounds: dict | None = None):
    """取得 body placeholder,優先選擇 rotation == 0 的(Bug 3 修正)"""
    from ._safe_shape import get_or_create_body

    return get_or_create_body(slide, slide_bounds)
