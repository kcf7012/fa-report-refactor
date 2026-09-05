"""問題描述與定義 Improver

針對「問題描述與定義」維度(權重 15%)。
改善內容:
- 失效現象 vs 失效模式 對照表
- 問題範圍量化(失效率、影響數量)
- 客戶影響評估
"""

from __future__ import annotations

from pptx import Presentation

from ..domain.evaluation import DimensionScore, EvaluationResult
from ..layout.selector import find_content_layout
from ..templates.loader import TemplateLoader
from ..visuals import (
    ELAN_BLUE,
    ELAN_ORANGE,
    ELAN_RED,
    ChecklistGenerator,
    ComparisonTableGenerator,
)
from ._safe_shape import BODY_SAFE_LEFT_INCH
from ._template_helper import resolve_template


def add_problem_definition_slide(
    prs: Presentation,
    evaluation: EvaluationResult,
    dimension: DimensionScore | None = None,
    template_loader: TemplateLoader | None = None,
    template_name: str = "problem_definition",
    slide_bounds: dict | None = None,
) -> None:
    """新增問題描述與定義投影片

    Args:
        prs: 簡報物件
        evaluation: 評估結果
        dimension: 該維度的評分(可選,用於針對性改善)
        template_loader: 樣板載入器(可選)
        template_name: 樣板名稱(預設 'problem_definition')
        slide_bounds: slide 尺寸(英寸),動態適應
    """
    from ._logging import log_action

    with log_action("add_problem_definition_slide"):
        layout = find_content_layout(prs)
        slide = prs.slides.add_slide(layout)

        # === 動態座標 ===
        sw = slide_bounds["width_inch"] if slide_bounds else 10.0
        # v3.1.4:與 _safe_shape safe_textbox fallback 對齊(margin=1.0,確保小 slide 不會太擠)
        margin = BODY_SAFE_LEFT_INCH
        if margin < 0.5:
            margin = 0.5
        content_w = sw - 2 * margin

        # 標題 — 優先從樣板讀取
        title = _get_or_create_title(slide, slide_bounds)
        if template_loader is not None:
            try:
                template = resolve_template(template_loader, template_name)
                title.text_frame.text = template.title
            except KeyError:
                title.text_frame.text = "問題描述與失效定義"
        else:
            title.text_frame.text = "問題描述與失效定義"

        # 1. 失效現象 vs 失效模式對照表
        _add_phenomenon_vs_mode_table(slide, content_w)

        # 2. 問題範圍量化檢查清單
        _add_quantification_checklist(slide, content_w)


def _add_phenomenon_vs_mode_table(slide, content_w: float = 9.0) -> None:
    """加入失效現象 vs 失效模式對照表"""
    gen = ComparisonTableGenerator(
        slide,
        left=BODY_SAFE_LEFT_INCH,  # P4:body 安全左界(有量測依據)
        top=1.4,
        width=content_w,
        height=2.0,
    )
    gen.generate(
        {
            "headers": ["項目", "目前報告常見缺失", "建議補充內容"],
            "rows": [
                [
                    "失效現象 (Phenomenon)",
                    "僅描述「通訊失敗」等表面現象",
                    "完整描述:電壓/電流/波形/時序/外觀",
                ],
                [
                    "失效模式 (Failure Mode)",
                    "未明確定義失效機制",
                    "對應失效機制:開路/短路/漏電/漂移/功能失效",
                ],
                [
                    "失效位置",
                    "只說「IC 損壞」",
                    "精確位置:腳位/區塊/層次(晶圓/封裝/PCB)",
                ],
            ],
        }
    )


def _add_quantification_checklist(slide, content_w: float = 9.0) -> None:
    """加入問題範圍量化檢查清單"""
    gen = ChecklistGenerator(
        slide,
        left=BODY_SAFE_LEFT_INCH,  # P4:body 安全左界(有量測依據)
        top=3.7,
        width=content_w,
        height=3.0,
    )
    gen.generate(
        [
            {
                "text": "失效率(PPM 或 %):本次失效佔總出貨量的比例",
                "checked": False,
                "color": ELAN_RED,
            },
            {
                "text": "影響產品數量:目前庫存/在製品/已出貨的受影響數量",
                "checked": False,
                "color": ELAN_ORANGE,
            },
            {
                "text": "影響時間範圍:失效首次發生日期 → 目前(持續中/已結案)",
                "checked": False,
                "color": ELAN_ORANGE,
            },
            {
                "text": "客戶影響評估:客戶端失效比例、生產線停線時間、退貨/召回成本",
                "checked": False,
                "color": ELAN_RED,
            },
            {
                "text": "失效嚴重性分級:Critical(安全)/Major(功能)/Minor(性能降級)",
                "checked": False,
                "color": ELAN_RED,
            },
            {
                "text": "批量 vs 個案:是否為批次性失效(同批多顆)或個案(單顆)",
                "checked": False,
                "color": ELAN_BLUE,
            },
        ]
    )


def _get_or_create_title(slide, slide_bounds: dict | None = None):
    """取得真實的 title placeholder(Bug 2 + Bug 3 修正)"""
    from ._safe_shape import get_or_create_title

    return get_or_create_title(slide, slide_bounds)
