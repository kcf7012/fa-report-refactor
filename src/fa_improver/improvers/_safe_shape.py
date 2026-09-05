"""建立 shape 的共用 helper(防止 v3.1.1 的殘留版面渲染問題)

修正的 bug:
- Bug 3:textbox 被旋轉 90°(260811 slides 1/3/4/5/6)
- Bug 4:底部 placeholder 殘留「按一下即可新增文字」(N160JCN 多張)
- Bug 2:title placeholder 找不到正確的 title(N160JCN / MS 多張)

設計理由:把所有 shape 建立邏輯統一到 helper,
避免未來 improver 再用 raw `slide.shapes.add_textbox()` 而忘記設定 rotation。
"""

from __future__ import annotations

from pptx.util import Inches, Pt


def safe_textbox(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    text: str | None = None,
    font_size: int | None = None,
    font_bold: bool = False,
    word_wrap: bool = True,
    font_color_rgb: tuple[int, int, int] | None = None,
):
    """建立一個不會旋轉、不會變直式的 textbox

    Args:
        slide: pptx slide
        left/top/width/height: 英寸(float)
        text: 初始文字(可選)
        font_size: 字體大小 pt(可選)
        font_bold: 是否粗體
        word_wrap: 是否自動換行(預設 True)
        font_color_rgb: (R, G, B) tuple(可選)

    Returns:
        textbox shape
    """
    from pptx.dml.color import RGBColor

    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    # 關鍵:防止繼承母片的旋轉屬性(260811 bug)
    tb.rotation = 0

    tf = tb.text_frame
    tf.word_wrap = word_wrap
    # 關鍵:關掉 autofit,防止 textbox 寬度不足時自動縮小變直式
    tf.auto_size = None
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)

    if text is not None:
        p = tf.paragraphs[0]
        p.text = text
        if font_size is not None:
            p.font.size = Pt(font_size)
        p.font.bold = font_bold
        if font_color_rgb is not None:
            p.font.color.rgb = RGBColor(*font_color_rgb)

    return tb


def clean_unused_placeholders(slide, *, also_remove: bool = True) -> int:
    """清除未使用的 placeholder(預設文字或空)

    Args:
        slide: pptx slide
        also_remove: 是否從 slide 移除 placeholder(預設 True)
        also_clear_text: 已停用(被 also_remove 取代)

    Returns:
        清除的 placeholder 數量

    Note:
        只清空 placeholder 不夠 — LibreOffice 會 fallback 顯示
        layout 的預設文字(例如「按一下以編輯母片文字樣式」)。
        需要「從 slide 移除 placeholder」才能讓它們在 LibreOffice
        渲染時不顯示出來。

        _safe_shape_v2.py 的修正:從 slide 移除整個 placeholder 元素。
    """
    residual_markers = ("按一下", "Click to add", "Click here to add", "<click")
    cleared = 0
    # 用 list() 複製,避免迭代時改動
    for shape in list(slide.placeholders):
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        is_residual = not text or any(marker in text for marker in residual_markers)
        if is_residual and also_remove:
            # 從 slide 的 shape tree 中移除這個 placeholder
            sp = shape._element
            sp.getparent().remove(sp)
            cleared += 1
    return cleared


# ---------------------------------------------------------------------------
# 版面安全常數 —— 量測依據(2026-09-06,跨平台遷移 P4 第 1 步)
# ---------------------------------------------------------------------------
# 這些數字**先量測、後定值**。稽核警告過不要為了讓既有測試通過而反推數值,
# 所以下面先記錄量測結果,再由結果推常數,不是反過來。
#
# 量測腳本:`scripts/measure_master_decoration.py`
# 量測範圍:母片 + `find_content_layout()` 實際選中的那一個 layout
#           (所有 improver 都只用這一個 layout,其餘 layout 的裝飾永遠
#            不會出現在新投影片上,算進來只會讓數字虛高)
# 判定方式:非 placeholder 的 shape 即裝飾;只採計「靠左緣、右緣未過投影片
#           中線」這種**往右移就閃得開**的。滿版背景/橫幅(如 260811 母片
#           的 Group 39,0.00~9.68 in)移到哪都閃不開,不列入下限。
#
# 「可迴避左側裝飾」的最大右緣(in):
#
#   檔案                                    選中 layout        title 帶  body 帶
#   260811_Kobo_ZHT_RA6080_SPcomFailI       直排標題及文字        ——     1.12
#   MS_Meishan_ADO_445239_260716            2L - Topic          0.97      ——
#   N160JCN-EEK ... 260810                  Topic-Numbers       0.97      ——
#   synthetic_C_decoration                  Content w/ Caption  1.00      ——
#   (synthetic_A / synthetic_B 無左側裝飾)
#
#   title 帶(0.30~1.15 in)跨檔最大值:1.00 in
#     └─ synthetic_C 的 master/LeftTopDecoration(0.00~1.00,top 0.00~0.50)
#     └─ 真實客戶檔是 0.97:MS / N160JCN 的「群組」(0.54~0.97,top 0.00~0.94)
#   body 帶(1.50 in ~ 投影片底部 -0.5)跨檔最大值:1.12 in
#     └─ 260811 的 master/Picture 14(0.00~1.12,top 1.23~6.38)——
#        該 layout 沒有 showMasterSp="0",母片 shape 確實會被渲染出來
#
# 定值規則:**量測最大值 + 0.20 in 緩衝,向上取到 0.05 的倍數**。
#   title:1.00 + 0.20 = 1.20 → 1.20(與既有值相同 —— 量測獨立落回同一個數,
#          不是為了配合既有測試才這樣選)
#   body :1.12 + 0.20 = 1.32 → 1.35(既有的 `TITLE_SAFE_LEFT_INCH - 0.2` = 1.00
#          **不足**,比 260811 的 Picture 14 右緣還左 0.12 in,正文首字會壓在
#          母片裝飾上)
# ---------------------------------------------------------------------------

TITLE_SAFE_LEFT_INCH: float = 1.2
TITLE_SAFE_TOP_INCH: float = 0.3
TITLE_SAFE_HEIGHT_INCH: float = 0.85

# body 內容的安全左界。P4 之前是各 improver 各自寫 `TITLE_SAFE_LEFT_INCH - 0.2`
# (= 1.00 in),那個 `- 0.2` 沒有任何量測依據,而且**不足**:比 260811 母片
# Picture 14 的右緣(1.12 in)還左 0.12 in。改成獨立常數,值由上面的量測決定,
# 不再掛在 title 常數上跟著飄。
BODY_SAFE_LEFT_INCH: float = 1.35

# 原生 title placeholder 被往右移之後,寬度低於此值就放棄它、改用 safe_textbox。
# ⚠️ 這個數字**不是量測來的**,是可讀性下限的判斷值:2.0 in 在 24pt 下約 6 個
# 中文字,再窄下去 title 會被迫折行,比 fallback textbox(約 8.3 in 寬)更糟。
# 與上面那些有量測依據的常數分開標示,免得被誤讀成也有母片量測背書。
TITLE_MIN_WIDTH_INCH: float = 2.0

# body placeholder 最小可用高度(低於此值就 fallback 用 safe_textbox)
# 「Topic-Numbers」layout 的 body placeholder 只有 0.51 in 高,放不下 heading+bullets
BODY_MIN_HEIGHT_INCH: float = 1.0


def _ensure_title_left_safe(shape):
    """把命中的 title shape 往右移到安全左界,移不動就回 ``None``。

    P4 修正:`get_title_placeholder()` 的策略 1/2/3 原本一命中就直接
    `return ph`,**完全不看 `left`**,所以 `TITLE_SAFE_LEFT_INCH` 只在
    「找不到 placeholder」的 fallback 分支生效 —— 而那是少數情況。

    為什麼是「移動」而不是「回 None 改用 safe_textbox」:移動保留母片給
    placeholder 的字型 / 顏色樣式,符合本專案「母片保護是最高優先」的原則。
    只有移完寬度不夠時才降級。

    幾何寫在 **slide 層級**的 `<a:xfrm>`(python-pptx 會在該 slide 的
    `spPr` 建/改 `a:off`、`a:ext`),master 與 layout 的 XML 完全不動,
    所以 `tests/unit/test_master_protection.py` 不受影響。

    **不自己追繼承鏈**:python-pptx 的 `_InheritsDimensions._effective_value()`
    已經做完「本層有設就用本層,否則往 layout / master 取」,`shape.left`
    讀到的就是有效值。唯一要保留的是 `None` 判斷 —— layout 與 master 都
    沒有對應 placeholder 時 `_inherited_value()` 會回 `None`。

    寬度採「保留原本的右緣」而不是「補到投影片邊界」:右緣不動就不會擠到
    同一張 layout 上其他 placeholder(例如 `Content with Caption` 的
    `Content Placeholder 2` 就緊接在 title 右邊)。
    """
    left = shape.left
    if left is None:
        return shape  # 無座標可判定,維持原行為(不要假裝有防線)
    safe_left = Inches(TITLE_SAFE_LEFT_INCH)
    if left >= safe_left:
        return shape

    width = shape.width
    if width is None:
        return None  # 有 left 卻沒有 width,無法安全縮寬 → 交給 safe_textbox
    new_width = left + width - safe_left
    if new_width < Inches(TITLE_MIN_WIDTH_INCH):
        return None  # 移完太窄,fallback textbox 比較好

    shape.left = safe_left
    shape.width = new_width
    return shape


def get_title_placeholder(slide):
    """取得真實的 title placeholder

    修正 Bug 2:嚴格用 placeholder_format.idx == 0,
    而不是只看 shape.name 或 slide.shapes.title(可能在 MS / N160JCN
    母片設計中找到「按一下即可新增文字」副 placeholder)。

    同時檢查 layout name 是否含「直排」,若是則跳過 layout 的
    title placeholder(避免中文直排)。

    修正 v3.1.3:當 layout 沒有 idx=0 placeholder 且 placeholder 數量
    只有 1 個(如「Topic-Numbers」、「Topic」單 placeholder layout),
    跳過 layout placeholder(改用 safe_textbox fallback),避免
    title 與該 placeholder 重疊。

    Args:
        slide: pptx slide

    Returns:
        title placeholder shape,或 None
    """
    # 先檢查 layout(直排 layout 不應使用其 title placeholder)
    # ⚠️ 已知缺陷(第一輪稽核 backlog #2,P4 刻意不修,留給 P7):
    #    這是**脆弱的字串比對**,只涵蓋 zh-TW「直排」與 en「Vertical」。
    #    layout 名稱取決於建立該 pptx 的 PowerPoint UI 語言,所以
    #    zh-CN「竖排」、ja「縦書き」以及「垂直」「縱向」「Portrait」全會漏掉,
    #    Bug 3(90° 旋轉)可能重現而測不出來。
    #    **不要靠加關鍵字來修**:那是丟棄式工作,改讀 XML 屬性之後整個作廢,
    #    而且機制照樣是猜名稱,只是讓症狀變罕見。根治要讀 bodyPr 的 `vert`
    #    屬性 / placeholder 的 `orient`。P7 一併處理,見該節。
    layout_name = slide.slide_layout.name
    if "直排" in layout_name or "Vertical" in layout_name:
        # 跳過 layout placeholder,改用 safe_textbox fallback
        return None

    # P4:三條策略命中後都要過 _ensure_title_left_safe(),不能直接 return。
    # 回 None 代表「這個 placeholder 救不回來」,交給 get_or_create_title()
    # 用 safe_textbox 重建,不再往下試其他策略(往下找也是同一批 placeholder)。

    # 策略 1:用 placeholder_format.idx == 0(最嚴格)
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            return _ensure_title_left_safe(ph)
    # 策略 2:用 placeholder type 找 TITLE/CENTER_TITLE
    try:
        from pptx.enum.shapes import PP_PLACEHOLDER

        for ph in slide.placeholders:
            if ph.placeholder_format.type in (
                PP_PLACEHOLDER.TITLE,
                PP_PLACEHOLDER.CENTER_TITLE,
            ):
                return _ensure_title_left_safe(ph)
    except (AttributeError, ValueError):
        pass
    # 策略 3:fallback — 找名稱含 title 的 shape
    for shape in slide.shapes:
        if shape.has_text_frame and "title" in shape.name.lower():
            return _ensure_title_left_safe(shape)

    # P4:刪掉原本的 `if len(list(slide.placeholders)) <= 1: return None`
    # (第一輪稽核 backlog #1)。它與後面的 `return None` 結果完全相同,是死碼。
    #
    # 為什麼是刪掉、而不是把檢查移到策略 3 之前去「真正實作原始意圖」:
    # 原始意圖是「單一 placeholder 的 layout 要跳過,避免 title 與它重疊」,
    # 但 improver 建立的新投影片一律來自 `find_content_layout()`,而它明文
    # 要求 `placeholder_count >= 2`(見 layout/selector.py),所以「只有 1 個
    # placeholder」在這條路上到不了。把死碼改成會生效的分支,等於為一個
    # 不存在的情境新增行為變更。
    #
    # 該意圖底下真正的風險是「title 與 body 解析到同一個 shape」
    # (placeholder 型別為 TITLE 但 idx != 0 時,get_body_placeholder 的
    # `idx != 0` 條件會挑中同一個)。那與 placeholder 數量無關,現有六份
    # fixture 都到不了(PowerPoint 產生的 title 一律 idx=0),已記入 handoff
    # 待辦,不在 P4 動它。
    return None


def get_body_placeholder(slide):
    """取得 body placeholder,優先選擇 rotation == 0 且 orient='horiz' 的(Bug 3 修正)

    260811 的某些 layout body placeholder 被預設為:
    - layout name 含「直排」(例如「直排標題及文字」)
    - orient='vert'(垂直中文排版)
    - rotation=90°

    若直接寫入文字,整頁會變直式。

    修正策略:
    1. 先檢查 layout name 是否含「直排」→ 若是,直接跳過 layout placeholder
    2. 否則用 layout placeholder,但強制 orient='horiz' 與 rotation=0
    3. v3.1.3 新增:若 layout placeholder 高度 < BODY_MIN_HEIGHT_INCH(1.0 in),
       視為「太矮裝不下 heading+bullets」,fallback 讓 get_or_create_body
       用 safe_textbox 重新建立。
    """
    # 先檢查 layout
    # ⚠️ 與 get_title_placeholder() 完全相同的脆弱字串比對,同樣只涵蓋
    #    zh-TW / en,zh-CN「竖排」與 ja「縦書き」會漏掉。兩處要一起改,
    #    P4 刻意不動,留給 P7(根治方式:讀 bodyPr 的 `vert` 屬性)。
    layout_name = slide.slide_layout.name
    if "直排" in layout_name or "Vertical" in layout_name:
        # 跳過 layout placeholder
        return None

    for shape in slide.placeholders:
        if shape.placeholder_format.idx != 0:
            # v3.1.3 新增:若 placeholder 高度 < BODY_MIN_HEIGHT_INCH,
            # 跳過 — 改用 safe_textbox fallback(避免內容溢出)
            if shape.height is not None:
                height_inch = shape.height / 914400
                if height_inch < BODY_MIN_HEIGHT_INCH:
                    return None
            # 關鍵修正:改為水平 orient(即使原本是 vert)
            import contextlib

            with contextlib.suppress(AttributeError, KeyError):
                ph_elem = shape._element.find(
                    ".//{http://schemas.openxmlformats.org/presentationml/2006/main}ph"
                )
                if ph_elem is not None:
                    ph_elem.set("orient", "horiz")
            # 同時確保 rotation = 0
            with contextlib.suppress(AttributeError):
                shape.rotation = 0
            return shape
    return None


def get_or_create_title(slide, slide_bounds=None):
    """取得 title placeholder,若無則建立新 textbox(Bug 2 + Bug 3 修正)

    v3.1.3 修正:fallback 的 safe_textbox 從 left=0.5 改成 left=1.2 in,
    避免被母片左上裝飾(深藍直條 + 淺藍色塊在 x=0.54-0.97)擋住 title 第一個字。
    height 從 1.0 改成 0.85,更緊湊。
    """
    ph = get_title_placeholder(slide)
    if ph is not None:
        return ph
    sw = slide_bounds["width_inch"] if slide_bounds else 10.0
    # v3.1.3:左邊界用 TITLE_SAFE_LEFT_INCH(預設 1.2)而非 0.5
    left = TITLE_SAFE_LEFT_INCH
    # 確保 left 不超出 slide 寬度
    if left >= sw - 2.0:
        left = max(0.5, sw - 9.5)
    width = sw - left - 0.5
    return safe_textbox(
        slide,
        left=left,
        top=TITLE_SAFE_TOP_INCH,
        width=width,
        height=TITLE_SAFE_HEIGHT_INCH,
    )


def get_or_create_body(slide, slide_bounds=None):
    """取得 body placeholder,若無則建立新 textbox(Bug 3 修正)

    v3.1.3 修正:當 layout 的 body placeholder 高度 < BODY_MIN_HEIGHT_INCH
    (如「Topic-Numbers」layout 的 height=0.51)時,get_body_placeholder()
    已 return None,所以這裡會 fallback 建立新的 textbox,確保有足夠空間
    容納 heading + 多個 bullets,避免與 title 區重疊。
    """
    ph = get_body_placeholder(slide)
    if ph is not None:
        return ph

    sw = slide_bounds["width_inch"] if slide_bounds else 10.0
    sh = slide_bounds["height_inch"] if slide_bounds else 7.5
    # P4:改用有量測依據的 BODY_SAFE_LEFT_INCH。
    # (舊註解寫「左邊界比 title 多 0.1 in,確保對齊」,但程式其實是
    #  `TITLE_SAFE_LEFT_INCH - 0.2`,比 title 少 0.2 in —— 註解與程式碼
    #  從一開始就對不上,不要沿用那個說法。)
    margin = BODY_SAFE_LEFT_INCH
    if margin < 0.5:
        margin = 0.5
    return safe_textbox(
        slide,
        left=margin,
        top=1.5,
        width=sw - margin - 0.5,
        height=sh - 2.0,
    )
