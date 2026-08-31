"""新增改善對策投影片

從 TemplateLoader 載入 'prevention_overview' 樣板取得標題與 placeholder items。
向後相容:若不傳 loader,使用預設載入器。

視覺元素:使用 TimelineGenerator 呈現改善時程。
"""

from __future__ import annotations

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from ..domain.evaluation import EvaluationResult
from ..domain.suggestion import Improvement
from ..layout.selector import find_content_layout
from ..templates.loader import TemplateLoader
from ..visuals import ELAN_BLUE, ELAN_GREEN, ELAN_ORANGE, TimelineGenerator
from ._template_helper import get_resolved_placeholders, resolve_template


def add_prevention_measures_slide(
    prs: Presentation,
    evaluation: EvaluationResult,
    improvements: list[Improvement],
    template_loader: TemplateLoader | None = None,
    template_name: str = "prevention_overview",
    slide_bounds: dict | None = None,
) -> None:
    """新增長期預防措施與改善對策投影片

    Args:
        prs: 簡報物件
        evaluation: 評估結果
        improvements: 改進建議清單
        template_loader: 樣板載入器(可選)
        template_name: 樣板名稱(預設 'prevention_overview')
        slide_bounds: slide 尺寸(英寸),動態適應
    """
    from ._logging import log_action
    with log_action("add_prevention_measures_slide"):
        # === 動態座標 ===
        sw = slide_bounds["width_inch"] if slide_bounds else 10.0
        margin = 0.5
        content_w = sw - 2 * margin

        # 載入樣板
        template = resolve_template(template_loader, template_name)

        layout = find_content_layout(prs)
        slide = prs.slides.add_slide(layout)

        # 標題(從樣板)
        title = _get_or_create_title(slide, slide_bounds)
        title.text_frame.text = template.title

        body = _get_or_create_body(slide, slide_bounds)
        tf = body.text_frame
        tf.clear()

    # section 0:擬議改善對策項目
    section0 = template.sections[0] if template.sections else None
    heading0 = section0.heading if section0 else "擬議改善對策項目"

    p = tf.paragraphs[0]
    p.text = f"{heading0}:"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(0, 112, 192)

    # 從 improvements 抽取建議,限制在 max_bullets
    if improvements:
        max_bullets = section0.max_bullets if section0 else 3
        for imp in improvements[:max_bullets]:
            p = tf.add_paragraph()
            p.text = imp.suggestion
            p.font.size = Pt(14)

    # section 1:標準化與監測計畫
    section1 = template.sections[1] if len(template.sections) > 1 else None
    if section1:
        p = tf.add_paragraph()
        p.text = f"\n[{section1.heading}]"
        p.font.bold = True

        # 優先用樣板的 placeholder_items
        placeholders = get_resolved_placeholders(template, section_index=1)
        if placeholders:
            for item in placeholders:
                p = tf.add_paragraph()
                p.text = f"• {item}"
                p.font.size = Pt(12)
        else:
            # fallback
            for item in [
                "建立入料檢驗 (IQC) SOP 與測試閾值",
                "導入自動化監測設備於生產線",
                "將此案例納入知識管理資料庫以利後續追蹤",
            ]:
                p = tf.add_paragraph()
                p.text = f"• {item}"
                p.font.size = Pt(12)

    # 加入 TimelineGenerator 視覺化改善時程
    _add_prevention_timeline(slide, improvements, content_w)


def add_iqc_standard_slide(
    prs: Presentation,
    evaluation: EvaluationResult,
    template_loader: TemplateLoader | None = None,
    slide_bounds: dict | None = None,
) -> None:
    """新增 IQC 標準投影片(改善對策維度的第二張)

    內容:IQC 抽驗比例、抽驗項目、允收標準(AQL)、不合格處置流程。

    Args:
        prs: 簡報物件
        evaluation: 評估結果(可選)
        template_loader: 樣板載入器(可選)
        slide_bounds: slide 尺寸(英寸),由 orchestrator 傳入
    """
    from ._logging import log_action
    with log_action("add_iqc_standard_slide"):
        _add_prevention_subtype_slide(
            prs,
            title="IQC 入料檢驗標準",
            table_headers=["檢驗項目", "抽驗比例", "允收標準 (AQL)", "不合格處置"],
            table_rows=[
                ["外觀檢查", "AQL 1.0", "MIL-STD-105E Level II", "退貨重工"],
                ["尺寸量測", "5%", "±0.05 mm", "隔離/特採"],
                ["電性測試", "100%", "全數測試", "退貨/報廢"],
                ["X-ray 檢查", "10%", "無空洞、裂痕", "重工"],
                ["DPA 分析", "1 pcs/批", "結構完整", "退貨"],
            ],
            template_loader=template_loader,
            slide_bounds=slide_bounds,
        )


def add_monitoring_km_slide(
    prs: Presentation,
    evaluation: EvaluationResult,
    template_loader: TemplateLoader | None = None,
    slide_bounds: dict | None = None,
) -> None:
    """新增監測 KM(Knowledge Management)投影片(改善對策維度的第三張)

    內容:KMS 登錄欄位、定期 review 機制、跨部門分享管道。

    Args:
        prs: 簡報物件
        evaluation: 評估結果(可選)
        template_loader: 樣板載入器(可選)
        slide_bounds: slide 尺寸(英寸),由 orchestrator 傳入
    """
    from ._logging import log_action
    with log_action("add_monitoring_km_slide"):
        _add_prevention_subtype_slide(
            prs,
            title="監測與知識管理 (KM)",
            table_headers=["項目", "頻率", "負責單位", "產出文件"],
            table_rows=[
                ["KMS 登錄", "每案結案", "FAE", "KM 文件編號"],
                ["失效資料庫更新", "每月", "QRA", "失效資料庫月報"],
                ["跨部門分享會議", "每季", "PM", "會議紀錄"],
                ["同類型失效再發監測", "每月", "QRA", "失效趨勢分析"],
                ["SOP 改版提案", "每半年", "製程工程師", "SOP 改版單"],
            ],
            template_loader=template_loader,
            slide_bounds=slide_bounds,
        )


def _add_prevention_subtype_slide(
    prs: Presentation,
    title: str,
    table_headers: list[str],
    table_rows: list[list[str]],
    template_loader: TemplateLoader | None = None,
    slide_bounds: dict | None = None,
) -> None:
    """共用 helper:建立「標題 + 對照表 + 底部說明」的投影片

    Args:
        prs: 簡報物件
        title: 投影片標題
        table_headers: 對照表 headers
        table_rows: 對照表 rows
        template_loader: 樣板載入器(向後相容)
        slide_bounds: slide 尺寸(英寸)
    """
    from pptx.util import Inches
    from pptx.util import Pt as _Pt

    layout = find_content_layout(prs)
    slide = prs.slides.add_slide(layout)

    # 動態座標
    sw = slide_bounds["width_inch"] if slide_bounds else 10.0
    margin = 0.5
    content_w = sw - 2 * margin

    # Title
    title_shape = _get_or_create_title(slide)
    title_shape.text_frame.text = title

    # 對照表(從 header 算起)
    from ..visuals import ComparisonTableGenerator

    gen = ComparisonTableGenerator(
        slide,
        left=margin,
        top=1.4,
        width=content_w,
        height=4.5,
    )
    gen.generate({"headers": table_headers, "rows": table_rows})

    # 底部說明
    note_box = slide.shapes.add_textbox(
        Inches(margin), Inches(6.1), Inches(content_w), Inches(0.8)
    )
    tf = note_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "✓ 此頁為 v3.1.1 新增的標準化對照表,提供 IQC 抽驗比例與 KM 監測頻率"
    from pptx.dml.color import RGBColor
    p.font.color.rgb = RGBColor(0x80, 0x80, 0x80)
    p.font.size = _Pt(10)


def _add_prevention_timeline(
    slide, improvements: list[Improvement], content_w: float = 9.0
) -> None:
    """加入改善時程時間軸

    預設 3 個階段:短期(< 1 個月)、中期(1-3 個月)、長期(> 3 個月)
    """
    timeline_items = [
        {
            "label": "短期 (< 1 個月)",
            "timeframe": "D+0 ~ D+30",
            "color": ELAN_ORANGE,
            "detail": "建立 SOP 與緊急對策",
        },
        {
            "label": "中期 (1-3 個月)",
            "timeframe": "D+30 ~ D+90",
            "color": ELAN_BLUE,
            "detail": "導入監測設備與人員訓練",
        },
        {
            "label": "長期 (> 3 個月)",
            "timeframe": "D+90 ~ D+180",
            "color": ELAN_GREEN,
            "detail": "成效追蹤與知識庫建立",
        },
    ]

    timeline_gen = TimelineGenerator(
        slide,
        left=0.5,
        top=5.5,
        width=content_w,
        height=1.5,
    )
    timeline_gen.generate(timeline_items)


def _get_or_create_title(slide, slide_bounds: dict | None = None):
    if slide.shapes.title:
        return slide.shapes.title
    for shape in slide.shapes:
        if "title" in shape.name.lower():
            return shape
    sw = slide_bounds["width_inch"] if slide_bounds else 10.0
    margin = 0.5
    return slide.shapes.add_textbox(
        Inches(margin), Inches(0.3), Inches(sw - 2 * margin), Inches(1)
    )


def _get_or_create_body(slide, slide_bounds: dict | None = None):
    for shape in slide.placeholders:
        if shape.placeholder_format.idx != 0:
            return shape
    sw = slide_bounds["width_inch"] if slide_bounds else 10.0
    sh = slide_bounds["height_inch"] if slide_bounds else 7.5
    margin = 0.5
    return slide.shapes.add_textbox(
        Inches(margin), Inches(1.5), Inches(sw - 2 * margin), Inches(sh - 2.0)
    )
